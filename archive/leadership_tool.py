import streamlit as st
import os
import PyPDF2
import docx
import pandas as pd
from datetime import datetime, timedelta
import threading
import time
from src.llm import chat
from src.auth import JiraConfig
from src.jira_client import JiraClient
from src.agent_router import route, parse_assignee, wants_stories_only

# Import reportlab at module level
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# Import python-pptx at module level
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Set page config
st.set_page_config(page_title="Leadership Quality Assistant", page_icon="🧭", layout="wide")

# Add custom CSS for fixed sidebar scrolling
st.markdown("""
<style>
    /* Fixed sidebar scrolling */
    .css-1d391kg {
        overflow-y: auto !important;
        max-height: 100vh !important;
        position: fixed !important;
    }
    
    /* Ensure main content area adjusts */
    .main .block-container {
        padding-left: 1rem;
        padding-right: 1rem;
        max-width: 100%;
    }
    
    /* Custom scrollbar for sidebar */
    .css-1d391kg::-webkit-scrollbar {
        width: 6px;
    }
    
    .css-1d391kg::-webkit-scrollbar-track {
        background: #f1f1f1;
        border-radius: 3px;
    }
    
    .css-1d391kg::-webkit-scrollbar-thumb {
        background: #888;
        border-radius: 3px;
    }
    
    .css-1d391kg::-webkit-scrollbar-thumb:hover {
        background: #555;
    }
</style>
""", unsafe_allow_html=True)

# Set API key
os.environ["OPENAI_API_KEY"] = "your_openai_api_key_here"

# Create downloads folder and cleanup functions
def ensure_downloads_folder():
    """Ensure downloads folder exists"""
    downloads_folder = "downloads"
    if not os.path.exists(downloads_folder):
        os.makedirs(downloads_folder)
    return downloads_folder

def cleanup_old_files():
    """Clean up files older than 10 minutes"""
    downloads_folder = ensure_downloads_folder()
    current_time = datetime.now()
    cutoff_time = current_time - timedelta(minutes=10)
    
    try:
        for filename in os.listdir(downloads_folder):
            file_path = os.path.join(downloads_folder, filename)
            if os.path.isfile(file_path):
                file_time = datetime.fromtimestamp(os.path.getctime(file_path))
                if file_time < cutoff_time:
                    os.remove(file_path)
                    print(f"Cleaned up old file: {filename}")
    except Exception as e:
        print(f"Error during cleanup: {e}")

def schedule_cleanup():
    """Schedule cleanup to run every 5 minutes"""
    def cleanup_loop():
        while True:
            cleanup_old_files()
            time.sleep(300)  # 5 minutes
    
    # Start cleanup thread
    cleanup_thread = threading.Thread(target=cleanup_loop, daemon=True)
    cleanup_thread.start()

# Initialize cleanup
schedule_cleanup()

# Document processing functions
def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_docx(docx_file):
    """Extract text from DOCX file"""
    try:
        doc = docx.Document(docx_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_text_from_txt(txt_file):
    """Extract text from TXT file"""
    try:
        content = txt_file.read()
        if isinstance(content, bytes):
            return content.decode('utf-8').strip()
        else:
            return content.strip()
    except Exception as e:
        return f"Error reading TXT: {e}"

def detect_jira_content(document_text):
    """Detect if document contains Jira-related content"""
    jira_indicators = [
        'jira', 'sprint', 'story', 'bug', 'epic', 'task', 'backlog',
        'scrum', 'kanban', 'agile', 'project', 'issue', 'ticket',
        'assignee', 'reporter', 'priority', 'status', 'resolution',
        'PROJ-', 'DEV-', 'TEST-', 'BUG-', 'STORY-', 'EPIC-'
    ]
    
    text_lower = document_text.lower()
    matches = [indicator for indicator in jira_indicators if indicator in text_lower]
    
    # Extract potential Jira issue keys (e.g., PROJ-123)
    import re
    jira_keys = re.findall(r'[A-Z]+-\d+', document_text)
    
    return {
        'is_jira_related': len(matches) > 0 or len(jira_keys) > 0,
        'jira_indicators': matches,
        'jira_keys': jira_keys,
        'confidence': len(matches) + len(jira_keys)
    }

def process_document(uploaded_file):
    """Process uploaded document and extract text"""
    if uploaded_file is None:
        return ""
        
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_extension == 'docx':
        return extract_text_from_docx(uploaded_file)
    elif file_extension == 'txt':
        return extract_text_from_txt(uploaded_file)
    else:
        return f"Unsupported file format: {file_extension}"

# Jira functions
def get_current_sprint(jira_client, board_id):
    """Get current sprint from Jira"""
    try:
        return jira_client.get_current_sprint()
    except Exception as e:
        st.sidebar.warning(f"Could not get current sprint: {e}")
        return None

def search_jira_issues(jira_client, jql, max_results=50):
    """Search Jira issues"""
    try:
        print(f"🔍 Jira Client Config:")
        print(f"   Base URL: {jira_client.cfg.base_url}")
        print(f"   Email: {jira_client.cfg.email}")
        print(f"   Board ID: {jira_client.cfg.board_id}")
        print(f"   JQL Query: {jql}")
        
        result = jira_client.search(jql=jql, max_results=max_results)
        print(f"🔍 Raw Jira Response: {result}")
        return result
    except Exception as e:
        print(f"❌ Jira search error: {e}")
        st.sidebar.error(f"Jira search failed: {e}")
        return None

def jira_tool(query, jira_client, board_id):
    """Process Jira queries"""
    try:
        # Check if query contains specific Jira ticket (e.g., CCM-283)
        import re
        ticket_match = re.search(r'([A-Z]+-\d+)', query.upper())
        
        # Also check for follow-up queries about specific tickets
        description_keywords = ['description', 'details', 'full description', 'complete description', 'more details', 'tell me more']
        is_description_query = any(keyword in query.lower() for keyword in description_keywords)
        
        if ticket_match or (is_description_query and st.session_state.get("last_jira_ticket")):
            # Search for specific ticket
            ticket_key = ticket_match.group(1) if ticket_match else st.session_state.get("last_jira_ticket")
            jql = f"key = {ticket_key}"
            
            # Store the ticket for follow-up queries
            st.session_state.last_jira_ticket = ticket_key
            
            # Debug logging
            print(f"🔍 Searching for ticket: {ticket_key}")
            print(f"🔍 JQL: {jql}")
            
            results = search_jira_issues(jira_client, jql)
            
            # Debug logging
            print(f"🔍 Search results: {results}")
            
            if results and results.get("issues"):
                issues = []
                for issue in results["issues"]:
                    fields = issue.get("fields", {})
                    issues.append({
                        "key": issue.get("key"),
                        "summary": fields.get("summary", ""),
                        "status": fields.get("status", {}).get("name", ""),
                        "assignee": fields.get("assignee", {}).get("displayName", "Unassigned"),
                        "type": fields.get("issuetype", {}).get("name", ""),
                        "updated": fields.get("updated", ""),
                        "description": fields.get("description", ""),
                        "link": f"{jira_client.cfg.base_url}/browse/{issue.get('key')}"
                    })
                print(f"🔍 Found {len(issues)} issues")
                return {"tool": "JIRA", "jql": jql, "items": issues}
            else:
                print(f"🔍 No issues found for {ticket_key}")
                
                # Try alternative search approaches
                print(f"🔍 Trying alternative search approaches...")
                
                # Try searching in all projects
                alt_jql = f"key = {ticket_key} OR summary ~ \"{ticket_key}\""
                print(f"🔍 Alternative JQL: {alt_jql}")
                alt_results = search_jira_issues(jira_client, alt_jql)
                
                if alt_results and alt_results.get("issues"):
                    issues = []
                    for issue in alt_results["issues"]:
                        fields = issue.get("fields", {})
                        issues.append({
                            "key": issue.get("key"),
                            "summary": fields.get("summary", ""),
                            "status": fields.get("status", {}).get("name", ""),
                            "assignee": fields.get("assignee", {}).get("displayName", "Unassigned"),
                            "type": fields.get("issuetype", {}).get("name", ""),
                            "updated": fields.get("updated", ""),
                            "description": fields.get("description", ""),
                            "link": f"{jira_client.cfg.base_url}/browse/{issue.get('key')}"
                        })
                    print(f"🔍 Found {len(issues)} issues with alternative search")
                    return {"tool": "JIRA", "jql": alt_jql, "items": issues}
                
                return {"tool": "JIRA", "jql": jql, "items": [], "message": f"No ticket found with key {ticket_key}. Please check if the ticket exists and you have access to it."}
        
        # Parse assignee and issue type for general queries
        assignee = parse_assignee(query) or ""
        issue_type = "Story" if wants_stories_only(query) else None
        
        # Debug logging
        print(f"🔍 Parsed assignee: {assignee}")
        print(f"🔍 Issue type: {issue_type}")
        
        # Get current sprint
        sprint = get_current_sprint(jira_client, board_id)
        sprint_id = sprint.get("id") if sprint else None
        
        # Build JQL - only add sprint filter if user explicitly mentions sprint
        jql_parts = []
        
        # Check if user explicitly mentioned sprint in their query
        sprint_keywords = ['sprint', 'current sprint', 'this sprint', 'active sprint']
        user_mentioned_sprint = any(keyword in query.lower() for keyword in sprint_keywords)
        
        if assignee:
            jql_parts.append(f'assignee = "{assignee}"')
        
        # Only add sprint filter if user explicitly mentioned sprint
        if user_mentioned_sprint and sprint_id:
            jql_parts.append(f'sprint = {sprint_id}')
        elif not assignee and not user_mentioned_sprint:
            # If no assignee and no sprint mentioned, default to project only
            jql_parts.append("project = CCM")
        
        if issue_type:
            jql_parts.append(f'issuetype = "{issue_type}"')
        
        jql = " AND ".join(jql_parts) if jql_parts else "project = CCM"
        jql += " ORDER BY updated DESC"
        
        # Debug logging
        print(f"🔍 Jira Client Config:")
        print(f"   Base URL: {jira_client.cfg.base_url}")
        print(f"   Email: {jira_client.cfg.email}")
        print(f"   Board ID: {jira_client.cfg.board_id}")
        print(f"   JQL Query: {jql}")
        
        # Search issues
        results = search_jira_issues(jira_client, jql)
        
        if results and results.get("issues"):
            issues = []
            for issue in results["issues"][:10]:  # Limit to 10 issues
                fields = issue.get("fields", {})
                issues.append({
                    "key": issue.get("key"),
                    "summary": fields.get("summary", ""),
                    "status": fields.get("status", {}).get("name", ""),
                    "assignee": fields.get("assignee", {}).get("displayName", "Unassigned"),
                    "type": fields.get("issuetype", {}).get("name", ""),
                    "updated": fields.get("updated", ""),
                    "link": f"{jira_client.cfg.base_url}/browse/{issue.get('key')}"
                })
            return {"tool": "JIRA", "jql": jql, "items": issues}
        else:
            # If no results and we have an assignee, try a broader search without sprint constraint
            if assignee:
                print(f"🔍 No results in current sprint, trying broader search for {assignee}...")
                broader_jql = f'assignee = "{assignee}" ORDER BY updated DESC'
                print(f"🔍 Broader JQL: {broader_jql}")
                broader_results = search_jira_issues(jira_client, broader_jql)
                
                if broader_results and broader_results.get("issues"):
                    issues = []
                    for issue in broader_results["issues"][:10]:  # Limit to 10 issues
                        fields = issue.get("fields", {})
                        issues.append({
                            "key": issue.get("key"),
                            "summary": fields.get("summary", ""),
                            "status": fields.get("status", {}).get("name", ""),
                            "assignee": fields.get("assignee", {}).get("displayName", "Unassigned"),
                            "type": fields.get("issuetype", {}).get("name", ""),
                            "updated": fields.get("updated", ""),
                            "link": f"{jira_client.cfg.base_url}/browse/{issue.get('key')}"
                        })
                    print(f"🔍 Found {len(issues)} issues in broader search")
                    return {"tool": "JIRA", "jql": broader_jql, "items": issues}
            
            return {"tool": "JIRA", "jql": jql, "items": [], "message": "No issues found"}
            
    except Exception as e:
        return {"tool": "JIRA", "error": str(e)}

def export_to_excel(messages):
    """Export assistant responses to Excel with proper formatting"""
    try:
        # Filter only assistant responses
        assistant_messages = [msg for msg in messages if msg["role"] == "assistant"]
        
        if not assistant_messages:
            st.warning("No assistant responses to export")
            return None
        
        # Create Excel file with better naming
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = f"Leadership_Responses_{timestamp}.xlsx"
        downloads_folder = ensure_downloads_folder()
        filepath = os.path.join(downloads_folder, filename)
        
        # Import openpyxl for advanced formatting
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
        
        # Create workbook and worksheet
        wb = Workbook()
        ws = wb.active
        ws.title = "AI Responses"
        
        # Define styles
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        content_font = Font(size=11)
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        # Set headers
        ws['A1'] = "Response #"
        ws['B1'] = "Content"
        ws['C1'] = "Timestamp"
        
        # Apply header styling
        for col in ['A1', 'B1', 'C1']:
            ws[col].font = header_font
            ws[col].fill = header_fill
            ws[col].alignment = Alignment(horizontal='center', vertical='center')
            ws[col].border = thin_border
        
        # Process each response
        row = 2
        for i, message in enumerate(assistant_messages, 1):
            content = message["content"]
            
            # Clean and format content
            formatted_content = format_content_for_excel(content)
            
            # Set cell values
            ws[f'A{row}'] = i
            ws[f'B{row}'] = formatted_content
            ws[f'C{row}'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            
            # Apply styling
            ws[f'A{row}'].font = content_font
            ws[f'A{row}'].alignment = Alignment(horizontal='center', vertical='top')
            ws[f'A{row}'].border = thin_border
            
            ws[f'B{row}'].font = content_font
            ws[f'B{row}'].alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            ws[f'B{row}'].border = thin_border
            
            ws[f'C{row}'].font = content_font
            ws[f'C{row}'].alignment = Alignment(horizontal='center', vertical='top')
            ws[f'C{row}'].border = thin_border
            
            row += 1
            
            # Auto-adjust column widths
        ws.column_dimensions['A'].width = 12  # Response #
        ws.column_dimensions['B'].width = 80  # Content (wider for readability)
        ws.column_dimensions['C'].width = 20  # Timestamp
        
        # Set row heights for better readability
        for row_num in range(2, row):
            ws.row_dimensions[row_num].height = 60  # Allow for wrapped text
        
        # Save workbook
        wb.save(filepath)
        return filepath
    except Exception as e:
        st.error(f"Error creating Excel file: {e}")
        return None

def format_content_for_excel(content):
    """Format content to match UI display"""
    import re
    
    # Remove markdown formatting and convert to plain text with proper structure
    formatted = content
    
    # Convert bold text (**text** -> TEXT)
    formatted = re.sub(r'\*\*(.*?)\*\*', r'\1', formatted)
    
    # Convert bullet points (- -> •)
    formatted = re.sub(r'^- ', '• ', formatted, flags=re.MULTILINE)
    
    # Convert numbered lists (1. -> 1.)
    formatted = re.sub(r'^(\d+)\. ', r'\1. ', formatted, flags=re.MULTILINE)
    
    # Handle tables - convert to readable format
    if '|' in formatted and '\n' in formatted:
        lines = formatted.split('\n')
        table_lines = [line for line in lines if '|' in line and not line.strip().startswith('|')]
        if table_lines:
            # Convert table to readable format
            table_content = []
            for line in table_lines:
                if line.strip():
                    cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                    if cells:
                        table_content.append(' | '.join(cells))
            formatted = '\n'.join(table_content)
    
    # Clean up extra whitespace
    formatted = re.sub(r'\n\s*\n', '\n\n', formatted)  # Multiple newlines to double
    formatted = formatted.strip()
    
    return formatted

def export_to_pdf(messages):
    """Export assistant responses to PDF with proper formatting"""
    if not REPORTLAB_AVAILABLE:
        st.error("PDF export requires reportlab. Install with: pip install reportlab")
        return None
    
    try:
        # Filter only assistant responses
        assistant_messages = [msg for msg in messages if msg["role"] == "assistant"]
        
        if not assistant_messages:
            st.warning("No assistant responses to export")
            return None
        
        # Create PDF file with better naming
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = f"Leadership_Responses_{timestamp}.pdf"
        downloads_folder = ensure_downloads_folder()
        filepath = os.path.join(downloads_folder, filename)
        doc = SimpleDocTemplate(filepath, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
        
        # Styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1,  # Center
            textColor=colors.HexColor('#2c3e50')
        )
        
        response_header_style = ParagraphStyle(
            'ResponseHeader',
            parent=styles['Heading3'],
            fontSize=14,
            spaceAfter=10,
            textColor=colors.HexColor('#A23B72'),
            fontName='Helvetica-Bold'
        )
        
        content_style = ParagraphStyle(
            'ContentStyle',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=12,
            leading=14,
            leftIndent=0,
            rightIndent=0
        )
        
        # Content
        story = []
        
        # Title
        title = Paragraph("Leadership Quality Assistant - AI Responses", title_style)
        story.append(title)
        story.append(Spacer(1, 20))
        
        # Export info
        export_info = Paragraph(f"<b>Export Date:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal'])
        story.append(export_info)
        story.append(Spacer(1, 30))
        
        # Messages
        for i, message in enumerate(assistant_messages, 1):
            # Message header
            header = Paragraph(f"Response #{i}", response_header_style)
            story.append(header)
            
            # Message content
            content = message["content"]
            
            # Format content for PDF
            formatted_content = format_content_for_pdf(content)
            
            # Check if content contains tables (before HTML conversion)
            original_content = message["content"]
            if "|" in original_content and "\n" in original_content:
                lines = original_content.split('\n')
                table_lines = [line for line in lines if '|' in line and not line.strip().startswith('|')]
                if table_lines:
                    # Create a proper table
                    table_data = []
                    for line in table_lines:
                        if line.strip():
                            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                            if cells:
                                table_data.append(cells)
                    
                    if table_data:
                        table = Table(table_data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#366092')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 12),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('TOPPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8f9fa')),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black),
                            ('FONTSIZE', (0, 1), (-1, -1), 10)
                        ]))
                        story.append(table)
                        story.append(Spacer(1, 20))
                        continue
            
            # Regular content with proper formatting
            content_para = Paragraph(formatted_content, content_style)
            story.append(content_para)
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        return filepath
    except ImportError:
        st.error("PDF export requires reportlab. Install with: pip install reportlab")
        return None
    except Exception as e:
        st.error(f"Error creating PDF file: {e}")
        return None

def format_content_for_pdf(content):
    """Format content for PDF export to match UI display"""
    import re
    
    # Start with the original content
    formatted = content
    
    # Convert bold text (**text** -> <b>text</b>)
    formatted = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', formatted)
    
    # Convert bullet points (- -> •)
    formatted = re.sub(r'^- ', '• ', formatted, flags=re.MULTILINE)
    
    # Convert numbered lists (1. -> 1.)
    formatted = re.sub(r'^(\d+)\. ', r'\1. ', formatted, flags=re.MULTILINE)
    
    # Handle line breaks better for PDF
    formatted = re.sub(r'\n', '<br/>', formatted)
    
    # Clean up extra whitespace
    formatted = re.sub(r'<br/>\s*<br/>', '<br/><br/>', formatted)
    formatted = formatted.strip()
    
    return formatted

def export_to_word(messages):
    """Export assistant responses to Word document"""
    try:
        # Filter only assistant responses
        assistant_messages = [msg for msg in messages if msg["role"] == "assistant"]
        
        if not assistant_messages:
            st.warning("No assistant responses to export")
            return None
        
        # Create Word file with better naming
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = f"Leadership_Responses_{timestamp}.docx"
        downloads_folder = ensure_downloads_folder()
        filepath = os.path.join(downloads_folder, filename)
        
        # Import python-docx for Word document creation
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        # Create document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Leadership Quality Assistant - AI Responses', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add export info
        doc.add_paragraph(f'Export Date: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
        doc.add_paragraph('')  # Empty line
        
        # Process each response
        for i, message in enumerate(assistant_messages, 1):
            # Add response header
            doc.add_heading(f'Response #{i}', level=2)
            
            # Add content
            content = message["content"]
            
            # Format content for Word
            formatted_content = format_content_for_word(content)
            
            # Check if content contains tables
            if "|" in formatted_content and "\n" in formatted_content:
                lines = formatted_content.split('\n')
                table_lines = [line for line in lines if '|' in line and not line.strip().startswith('|')]
                if table_lines:
                    # Create a Word table
                    table_data = []
                    for line in table_lines:
                        if line.strip():
                            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                            if cells:
                                table_data.append(cells)
                    
                    if table_data:
                        # Create table
                        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                        table.style = 'Table Grid'
                        
                        # Add data to table
                        for i, row_data in enumerate(table_data):
                            for j, cell_data in enumerate(row_data):
                                table.cell(i, j).text = cell_data
                        
                        doc.add_paragraph('')  # Empty line after table
                        continue
            
            # Regular content
            doc.add_paragraph(formatted_content)
            doc.add_paragraph('')  # Empty line between responses
        
        # Save document
        doc.save(filepath)
        return filepath
    except ImportError:
        st.error("Word export requires python-docx. Install with: pip install python-docx")
        return None
    except Exception as e:
        st.error(f"Error creating Word file: {e}")
        return None

def format_content_for_word(content):
    """Format content for Word export to match UI display"""
    import re
    
    # Remove markdown formatting and convert to plain text with proper structure
    formatted = content
    
    # Convert bold text (**text** -> plain text)
    formatted = re.sub(r'\*\*(.*?)\*\*', r'\1', formatted)
    
    # Convert bullet points (- -> •)
    formatted = re.sub(r'^- ', '• ', formatted, flags=re.MULTILINE)
    
    # Convert numbered lists (1. -> 1.)
    formatted = re.sub(r'^(\d+)\. ', r'\1. ', formatted, flags=re.MULTILINE)
    
    # Handle tables - convert to readable format
    if '|' in formatted and '\n' in formatted:
        lines = formatted.split('\n')
        table_lines = [line for line in lines if '|' in line and not line.strip().startswith('|')]
        if table_lines:
            # Convert table to readable format
            table_content = []
            for line in table_lines:
                if line.strip():
                    cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                    if cells:
                        table_content.append(' | '.join(cells))
            formatted = '\n'.join(table_content)
    
    # Clean up extra whitespace
    formatted = re.sub(r'\n\s*\n', '\n\n', formatted)  # Multiple newlines to double
    formatted = formatted.strip()
    
    return formatted

def export_to_powerpoint(messages):
    """Export assistant responses to PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        st.error("PowerPoint export requires python-pptx. Install with: pip install python-pptx")
        return None
    
    try:
        # Filter only assistant responses
        assistant_messages = [msg for msg in messages if msg["role"] == "assistant"]
        
        if not assistant_messages:
            st.warning("No assistant responses to export")
            return None
        
        # Create PowerPoint file with better naming
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        filename = f"Leadership_Responses_{timestamp}.pptx"
        downloads_folder = ensure_downloads_folder()
        filepath = os.path.join(downloads_folder, filename)
        
        # Create presentation
        prs = Presentation()
        
        # Define colors
        title_color = RGBColor(44, 62, 80)  # Dark blue
        content_color = RGBColor(52, 73, 94)  # Medium blue
        accent_color = RGBColor(162, 59, 114)  # Purple
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Leadership Quality Assistant"
        subtitle.text = f"AI Responses Report\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Style title slide
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.color.rgb = content_color
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        
        # Content slides
        for i, message in enumerate(assistant_messages, 1):
            # Use content layout
            content_slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set slide title
            title_shape = slide.shapes.title
            title_shape.text = f"Response #{i}"
            
            # Style title
            title_shape.text_frame.paragraphs[0].font.color.rgb = accent_color
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Get content area
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Process content
            content = message["content"]
            formatted_content = format_content_for_powerpoint(content)
            
            # Add content to slide
            p = text_frame.paragraphs[0]
            p.text = formatted_content
            p.font.size = Pt(16)
            p.font.color.rgb = content_color
            p.alignment = PP_ALIGN.LEFT
            
            # Handle tables if present
            if "|" in content and "\n" in content:
                lines = content.split('\n')
                table_lines = [line for line in lines if '|' in line and not line.strip().startswith('|')]
                if table_lines and len(table_lines) > 1:
                    # Create a table slide
                    table_slide_layout = prs.slide_layouts[6]  # Blank layout
                    table_slide = prs.slides.add_slide(table_slide_layout)
                    
                    # Add title
                    title_shape = table_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    title_frame = title_shape.text_frame
                    title_frame.text = f"Response #{i} - Data Table"
                    title_frame.paragraphs[0].font.size = Pt(24)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = accent_color
                    
                    # Create table
                    table_data = []
                    for line in table_lines:
                        if line.strip():
                            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                            if cells:
                                table_data.append(cells)
                    
                    if table_data:
                        rows = len(table_data)
                        cols = len(table_data[0])
                        
                        # Add table
                        table_shape = table_slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4))
                        table = table_shape.table
                        
                        # Populate table
                        for i, row_data in enumerate(table_data):
                            for j, cell_data in enumerate(row_data):
                                if j < cols:  # Ensure we don't exceed column count
                                    cell = table.cell(i, j)
                                    cell.text = cell_data
                                    
                                    # Style header row
                                    if i == 0:
                                        cell.text_frame.paragraphs[0].font.bold = True
                                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = RGBColor(54, 96, 146)  # Blue header
                                    else:
                                        cell.text_frame.paragraphs[0].font.color.rgb = content_color
                                    
                                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Summary slide
        summary_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(summary_slide_layout)
        
        title_shape = summary_slide.shapes.title
        title_shape.text = "Summary"
        
        title_shape.text_frame.paragraphs[0].font.color.rgb = accent_color
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = summary_slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        summary_text = f"""• Total Responses: {len(assistant_messages)}
• Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
• Source: Leadership Quality Assistant
• Format: Professional Presentation

Thank you for using Leadership Quality Assistant!"""
        
        p = text_frame.paragraphs[0]
        p.text = summary_text
        p.font.size = Pt(18)
        p.font.color.rgb = content_color
        
        # Save presentation
        prs.save(filepath)
        return filepath
        
    except ImportError:
        st.error("PowerPoint export requires python-pptx. Install with: pip install python-pptx")
        return None
    except Exception as e:
        st.error(f"Error creating PowerPoint file: {e}")
        return None

def format_content_for_powerpoint(content):
    """Format content for PowerPoint export"""
    import re
    
    # Start with the original content
    formatted = content
    
    # Convert bold text (**text** -> plain text with emphasis)
    formatted = re.sub(r'\*\*(.*?)\*\*', r'\1', formatted)
    
    # Convert bullet points (- -> •)
    formatted = re.sub(r'^- ', '• ', formatted, flags=re.MULTILINE)
    
    # Convert numbered lists (1. -> 1.)
    formatted = re.sub(r'^(\d+)\. ', r'\1. ', formatted, flags=re.MULTILINE)
    
    # Remove table formatting for main content (tables get separate slides)
    if '|' in formatted and '\n' in formatted:
        lines = formatted.split('\n')
        non_table_lines = [line for line in lines if '|' not in line or line.strip().startswith('|')]
        formatted = '\n'.join(non_table_lines)
    
    # Clean up extra whitespace
    formatted = re.sub(r'\n\s*\n', '\n\n', formatted)
    formatted = formatted.strip()
    
    # Limit content length for slide readability (increased limit)
    if len(formatted) > 1000:
        formatted = formatted[:1000] + "..."
    
    return formatted

# Analytics and Dashboard Functions
def get_sprint_analytics(jira_client, board_id):
    """Get comprehensive sprint analytics"""
    try:
        # Get current sprint
        current_sprint = jira_client.get_current_sprint()
        if not current_sprint:
            return None
        
        sprint_id = current_sprint.get('id')
        sprint_name = current_sprint.get('name', 'Unknown')
        
        # Get sprint issues
        jql = f"sprint = {sprint_id}"
        sprint_issues = jira_client.search(jql=jql, max_results=100)
        
        if not sprint_issues.get('issues'):
            return {
                'sprint_name': sprint_name,
                'total_issues': 0,
                'completed_issues': 0,
                'in_progress_issues': 0,
                'todo_issues': 0,
                'completion_rate': 0,
                'team_performance': {},
                'bug_analysis': {},
                'velocity_data': []
            }
        
        issues = sprint_issues['issues']
        
        # Analyze issues by status
        status_counts = {}
        team_performance = {}
        bug_analysis = {'total_bugs': 0, 'critical_bugs': 0, 'resolved_bugs': 0}
        
        for issue in issues:
            fields = issue.get('fields', {})
            status = fields.get('status', {}).get('name', 'Unknown')
            assignee = fields.get('assignee', {}).get('displayName', 'Unassigned')
            issue_type = fields.get('issuetype', {}).get('name', 'Unknown')
            priority = fields.get('priority', {}).get('name', 'Unknown')
            
            # Count by status
            status_counts[status] = status_counts.get(status, 0) + 1
            
            # Team performance
            if assignee not in team_performance:
                team_performance[assignee] = {'total': 0, 'completed': 0, 'in_progress': 0}
            team_performance[assignee]['total'] += 1
            
            if status in ['Done', 'Closed', 'Resolved']:
                team_performance[assignee]['completed'] += 1
            elif status in ['In Progress', 'In Review']:
                team_performance[assignee]['in_progress'] += 1
            
            # Bug analysis
            if issue_type == 'Bug':
                bug_analysis['total_bugs'] += 1
                if priority in ['Critical', 'Highest']:
                    bug_analysis['critical_bugs'] += 1
                if status in ['Done', 'Closed', 'Resolved']:
                    bug_analysis['resolved_bugs'] += 1
        
        # Calculate completion rate
        total_issues = len(issues)
        completed_issues = status_counts.get('Done', 0) + status_counts.get('Closed', 0) + status_counts.get('Resolved', 0)
        completion_rate = (completed_issues / total_issues * 100) if total_issues > 0 else 0
        
        return {
            'sprint_name': sprint_name,
            'total_issues': total_issues,
            'completed_issues': completed_issues,
            'in_progress_issues': status_counts.get('In Progress', 0) + status_counts.get('In Review', 0),
            'todo_issues': status_counts.get('To Do', 0) + status_counts.get('Open', 0),
            'completion_rate': round(completion_rate, 1),
            'team_performance': team_performance,
            'bug_analysis': bug_analysis,
            'status_breakdown': status_counts,
            'velocity_data': []  # Will be populated with historical data
        }
    except Exception as e:
        st.error(f"Error getting sprint analytics: {e}")
        return None

def get_velocity_trends(jira_client, board_id, sprints_back=5):
    """Get velocity trends for the last N sprints"""
    try:
        # This would require historical sprint data
        # For now, return mock data structure
        return [
            {'sprint': 'Sprint 1', 'velocity': 23, 'completed': 18, 'planned': 25},
            {'sprint': 'Sprint 2', 'velocity': 27, 'completed': 22, 'planned': 28},
            {'sprint': 'Sprint 3', 'velocity': 25, 'completed': 20, 'planned': 26},
            {'sprint': 'Sprint 4', 'velocity': 30, 'completed': 25, 'planned': 32},
            {'sprint': 'Sprint 5', 'velocity': 28, 'completed': 23, 'planned': 30}
        ]
    except Exception as e:
        st.error(f"Error getting velocity trends: {e}")
        return []

def generate_qa_metrics(jira_client, board_id):
    """Generate QA-specific metrics"""
    try:
        # Get all bugs from current sprint
        current_sprint = jira_client.get_current_sprint()
        if not current_sprint:
            return None
        
        sprint_id = current_sprint.get('id')
        
        # Bug analysis
        bug_jql = f"sprint = {sprint_id} AND issuetype = Bug"
        bug_issues = jira_client.search(jql=bug_jql, max_results=50)
        
        # Test case analysis (if available)
        test_jql = f"sprint = {sprint_id} AND issuetype = Test"
        test_issues = jira_client.search(jql=test_jql, max_results=50)
        
        bugs = bug_issues.get('issues', [])
        tests = test_issues.get('issues', [])
        
        # Analyze bug metrics
        bug_metrics = {
            'total_bugs': len(bugs),
            'critical_bugs': 0,
            'high_priority_bugs': 0,
            'resolved_bugs': 0,
            'open_bugs': 0,
            'avg_resolution_time': 0,
            'bug_trend': 'stable'
        }
        
        for bug in bugs:
            fields = bug.get('fields', {})
            priority = fields.get('priority', {}).get('name', 'Unknown')
            status = fields.get('status', {}).get('name', 'Unknown')
            
            if priority in ['Critical', 'Highest']:
                bug_metrics['critical_bugs'] += 1
            elif priority == 'High':
                bug_metrics['high_priority_bugs'] += 1
            
            if status in ['Done', 'Closed', 'Resolved']:
                bug_metrics['resolved_bugs'] += 1
            else:
                bug_metrics['open_bugs'] += 1
        
        # Test metrics
        test_metrics = {
            'total_tests': len(tests),
            'passed_tests': 0,
            'failed_tests': 0,
            'test_coverage': 0
        }
        
        for test in tests:
            fields = test.get('fields', {})
            status = fields.get('status', {}).get('name', 'Unknown')
            
            if status == 'Passed':
                test_metrics['passed_tests'] += 1
            elif status == 'Failed':
                test_metrics['failed_tests'] += 1
        
        return {
            'bug_metrics': bug_metrics,
            'test_metrics': test_metrics,
            'quality_score': calculate_quality_score(bug_metrics, test_metrics)
        }
    except Exception as e:
        st.error(f"Error generating QA metrics: {e}")
        return None

def calculate_quality_score(bug_metrics, test_metrics):
    """Calculate overall quality score"""
    try:
        # Simple quality score calculation
        bug_score = 100 - (bug_metrics['critical_bugs'] * 20) - (bug_metrics['high_priority_bugs'] * 10)
        test_score = (test_metrics['passed_tests'] / max(test_metrics['total_tests'], 1)) * 100
        
        quality_score = (bug_score + test_score) / 2
        return max(0, min(100, quality_score))
    except:
        return 0

# Main Title - Generic Integration Platform
st.markdown("""
<div style="display: flex; align-items: center; justify-content: center; margin: 20px 0; padding: 30px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 20px; box-shadow: 0 8px 32px rgba(0,0,0,0.1);">
    <div style="text-align: center;">
        <h1 style="margin: 0; color: white; font-size: 3rem; font-weight: bold; text-shadow: 2px 2px 4px rgba(0,0,0,0.3);">🔗 Leadership Quality Assistant</h1>
        <p style="margin: 10px 0 0 0; color: rgba(255,255,255,0.9); font-size: 1.3rem; font-weight: 500;">Connect • Analyze • Lead</p>
        <p style="margin: 5px 0 0 0; color: rgba(255,255,255,0.7); font-size: 1rem;">Powered by TAO DIGITAL</p>
    </div>
</div>
""", unsafe_allow_html=True)

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "document_text" not in st.session_state:
    st.session_state.document_text = ""
if "last_uploaded_file" not in st.session_state:
    st.session_state.last_uploaded_file = None
if "clear_clicked" not in st.session_state:
    st.session_state.clear_clicked = False
if "show_export_buttons" not in st.session_state:
    st.session_state.show_export_buttons = False
if "export_downloaded" not in st.session_state:
    st.session_state.export_downloaded = False
if "current_export_file" not in st.session_state:
    st.session_state.current_export_file = None
if "export_file_created" not in st.session_state:
    st.session_state.export_file_created = False
if "last_jira_ticket" not in st.session_state:
    st.session_state.last_jira_ticket = None

# Main Content Area - Centered Layout
col1, col2, col3 = st.columns([1, 2, 1])

with col2:
    # Welcome Section
    st.markdown("""
    <div style="text-align: center; margin: 30px 0; padding: 25px; background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); border-radius: 15px; box-shadow: 0 4px 15px rgba(0,0,0,0.1);">
        <h2 style="color: #2c3e50; margin-bottom: 15px;">🚀 Welcome to Integration Hub</h2>
        <p style="color: #6c757d; font-size: 1.1rem; margin-bottom: 20px;">Connect your tools, analyze your data, and get intelligent insights</p>
        <div style="display: flex; justify-content: center; gap: 20px; flex-wrap: wrap;">
            <div style="background: white; padding: 15px; border-radius: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); min-width: 120px;">
                <div style="font-size: 2rem; margin-bottom: 5px;">📊</div>
                <div style="font-weight: bold; color: #2c3e50;">Analytics</div>
            </div>
            <div style="background: white; padding: 15px; border-radius: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); min-width: 120px;">
                <div style="font-size: 2rem; margin-bottom: 5px;">📄</div>
                <div style="font-weight: bold; color: #2c3e50;">Documents</div>
            </div>
            <div style="background: white; padding: 15px; border-radius: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); min-width: 120px;">
                <div style="font-size: 2rem; margin-bottom: 5px;">💬</div>
                <div style="font-weight: bold; color: #2c3e50;">Chat</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Document Upload Section
    st.header("📄 Document Analysis")
    uploaded_file = st.file_uploader("Upload any document for intelligent analysis", type=['pdf', 'docx', 'txt'], help="Upload documents to get AI-powered insights and analysis")

# Process uploaded file
if uploaded_file is not None and not st.session_state.clear_clicked:
    # Check if this is a new file upload
    if uploaded_file != st.session_state.last_uploaded_file:
        st.session_state.last_uploaded_file = uploaded_file
        try:
            # Use proper document processing
            document_text = process_document(uploaded_file)
            if document_text and not document_text.startswith("Error"):
                st.session_state.document_text = document_text
                
                # Detect Jira content in the document
                jira_analysis = detect_jira_content(document_text)
                
                # Show temporary success message
                st.toast(f"✅ Document '{uploaded_file.name}' uploaded successfully!", icon="✅")
                
                # If Jira content detected and Jira is configured, fetch relevant info
                if jira_analysis['is_jira_related'] and st.session_state.get("jira_configured"):
                    st.warning("🔍 Jira-related content detected in document!")
                    st.info(f"Found indicators: {', '.join(jira_analysis['jira_indicators'])}")
                    if jira_analysis['jira_keys']:
                        st.info(f"Found Jira keys: {', '.join(jira_analysis['jira_keys'])}")
                    
                    # Automatically fetch relevant Jira information
                    try:
                        jira_client = st.session_state.get("jira_client")
                        board_id = st.session_state.get("jira_board_id")
                        
                        if jira_client and board_id:
                            # Get current sprint info
                            sprint = jira_client.get_current_sprint()
                            if sprint:
                                st.success(f"🏃 Current Sprint: {sprint.get('name', 'Unknown')}")
                            
                            # Search for issues mentioned in document
                            if jira_analysis['jira_keys']:
                                for key in jira_analysis['jira_keys'][:3]:  # Limit to 3 keys
                                    try:
                                        issue = jira_client.get_issue(key)
                                        if issue:
                                            fields = issue.get('fields', {})
                                            st.info(f"📋 {key}: {fields.get('summary', 'No summary')} - {fields.get('status', {}).get('name', 'Unknown status')}")
                                    except:
                                        pass
                    except Exception as e:
                        st.error(f"Could not fetch Jira information: {e}")
                
                elif jira_analysis['is_jira_related'] and not st.session_state.get("jira_configured"):
                    # Silently ignore Jira content detection when Jira is not configured
                    pass
            else:
                st.error(f"❌ {document_text}")
        except Exception as e:
            st.error(f"Error processing file: {e}")

# Reset clear flag after processing
if st.session_state.clear_clicked:
    st.session_state.clear_clicked = False

    # Chat Interface Section
    st.markdown("---")
    st.header("💬 Intelligent Chat")

# Display chat messages
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.write(message["content"])

# Chat input
    if prompt := st.chat_input("Ask me anything about your data, documents, or integrations..."):
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        # Display user message
        with st.chat_message("user"):
            st.write(prompt)
        
        # Get AI response
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                try:
                    # Check if it's actually a general question that shouldn't be routed to Jira
                    # Only exclude truly general questions, not work-related ones
                    general_keywords = ['weather', 'joke', 'time', 'date', 'news', 'music', 'movie', 'game', 'recipe', 'travel', 'shopping', 'politics', 'celebrity', 'hello', 'hi', 'how are you']
                    is_general_question = any(keyword in prompt.lower() for keyword in general_keywords)
                    
                    # Check for Jira-specific keywords that should always go to Jira when configured
                    jira_specific_keywords = ['ccm-', 'jira', 'ticket', 'bug', 'story', 'sprint', 'backlog', 'issue', 'epic', 'board', 'velocity', 'assignee', 'status', 'project', 'team', 'analytics', 'dashboard', 'metrics', 'performance', 'trends', 'report', 'summary', 'progress', 'completion', 'burndown', 'capacity', 'productivity']
                    is_jira_related = any(keyword in prompt.lower() for keyword in jira_specific_keywords)
                    
                    # Check for analytics queries
                    analytics_keywords = ['analytics', 'dashboard', 'metrics', 'velocity', 'performance', 'trends', 'report', 'summary', 'status', 'progress', 'completion', 'team performance', 'bug analysis', 'quality score', 'sprint health', 'burndown', 'capacity', 'productivity']
                    is_analytics_query = any(keyword in prompt.lower() for keyword in analytics_keywords)
                    
                    # Check for Excel format requests
                    excel_keywords = ['excel', 'spreadsheet', 'table format', 'tabular', 'csv']
                    is_excel_request = any(keyword in prompt.lower() for keyword in excel_keywords)
                    
                    # Check for PowerPoint requests
                    powerpoint_keywords = ['powerpoint', 'ppt', 'presentation', 'slides', 'create ppt', 'make ppt', 'slide deck']
                    is_powerpoint_request = any(keyword in prompt.lower() for keyword in powerpoint_keywords)
                    
                    # Determine if this is a Jira query - prioritize Jira when configured
                    if st.session_state.get("jira_configured") and is_jira_related:
                        decision = "JIRA"
                    elif not is_general_question:
                        decision = route(prompt)
                    else:
                        decision = "GENERAL"
                    
                    if st.session_state.document_text:
                        # Document Q&A - prioritize document analysis
                        if is_excel_request:
                            # Special handling for Excel format requests
                            full_prompt = f"""
                            Document: {st.session_state.document_text[:5000]}
                            
                            Question: {prompt}
                            
                            Format your response as a table with clear columns and rows. Use markdown table format with | separators.
                            Make it suitable for Excel export with proper headers and organized data.
                            """
                        elif is_powerpoint_request:
                            # Special handling for PowerPoint format requests
                            full_prompt = f"""
                            Document: {st.session_state.document_text[:5000]}
                            
                            Question: {prompt}
                            
                            Format your response for a PowerPoint presentation. Use clear headings, bullet points, and structured content.
                            Make it suitable for slide presentation with key points and professional formatting.
                            Include bold text for emphasis and organize information logically.
                            """
                        else:
                            # Regular document Q&A
                            full_prompt = f"""
                            Document: {st.session_state.document_text[:5000]}
                            
                            Question: {prompt}
                            
                            Provide a concise, professional answer based on the document.
                            """
                        
                        # Simplified response handling without threading
                        try:
                            response = chat([{"role": "user", "content": full_prompt}])
                        except Exception as e:
                            response = f"Error processing request: {str(e)}"
                    
                    elif decision == "JIRA" and st.session_state.get("jira_configured"):
                        # Jira query (only if not a general question)
                        jira_client = st.session_state.get("jira_client")
                        board_id = st.session_state.get("jira_board_id")
                        
                        if jira_client and board_id:
                            # Check if this is an analytics query
                            if is_analytics_query:
                                # Handle analytics queries
                                if 'sprint' in prompt.lower() and 'analytics' in prompt.lower():
                                    analytics = get_sprint_analytics(jira_client, board_id)
                                    if analytics:
                                        response = f"""📊 **Sprint Analytics Report**

**Sprint**: {analytics['sprint_name']}
**Total Issues**: {analytics['total_issues']}
**Completed**: {analytics['completed_issues']}
**Completion Rate**: {analytics['completion_rate']}%

**Team Performance:**
"""
                                        for member, perf in analytics['team_performance'].items():
                                            completion_rate = (perf['completed'] / perf['total'] * 100) if perf['total'] > 0 else 0
                                            response += f"• **{member}**: {perf['completed']}/{perf['total']} completed ({completion_rate:.1f}%)\n"
                                        
                                        if analytics['bug_analysis']['total_bugs'] > 0:
                                            response += f"\n**Bug Analysis:**\n"
                                            response += f"• Total Bugs: {analytics['bug_analysis']['total_bugs']}\n"
                                            response += f"• Critical Bugs: {analytics['bug_analysis']['critical_bugs']}\n"
                                            response += f"• Resolved: {analytics['bug_analysis']['resolved_bugs']}\n"
                                    else:
                                        response = "No sprint analytics data available."
                            
                            elif 'qa' in prompt.lower() or 'quality' in prompt.lower():
                                qa_metrics = generate_qa_metrics(jira_client, board_id)
                                if qa_metrics:
                                    response = f"""🔍 **QA Quality Report**

**Quality Score**: {qa_metrics['quality_score']:.1f}/100

**Bug Metrics:**
• Total Bugs: {qa_metrics['bug_metrics']['total_bugs']}
• Critical Bugs: {qa_metrics['bug_metrics']['critical_bugs']}
• High Priority: {qa_metrics['bug_metrics']['high_priority_bugs']}
• Resolved: {qa_metrics['bug_metrics']['resolved_bugs']}

**Test Metrics:**
• Total Tests: {qa_metrics['test_metrics']['total_tests']}
• Passed: {qa_metrics['test_metrics']['passed_tests']}
• Failed: {qa_metrics['test_metrics']['failed_tests']}
"""
                                    if qa_metrics['test_metrics']['total_tests'] > 0:
                                        success_rate = (qa_metrics['test_metrics']['passed_tests'] / qa_metrics['test_metrics']['total_tests'] * 100)
                                        response += f"• Success Rate: {success_rate:.1f}%\n"
                                else:
                                    response = "No QA metrics data available."
                            
                            elif 'velocity' in prompt.lower():
                                velocity_data = get_velocity_trends(jira_client, board_id)
                                if velocity_data:
                                    response = f"""📈 **Velocity Trends Report**

**Recent Sprint Performance:**
"""
                                    for sprint in velocity_data[-3:]:  # Last 3 sprints
                                        response += f"• **{sprint['sprint']}**: {sprint['velocity']} points (Completed: {sprint['completed']}, Planned: {sprint['planned']})\n"
                                    
                                    avg_velocity = sum(s['velocity'] for s in velocity_data) / len(velocity_data)
                                    response += f"\n**Average Velocity**: {avg_velocity:.1f} points"
                                else:
                                    response = "No velocity data available."
                            
                            else:
                                # General analytics query
                                response = """📊 **Analytics Dashboard Available**

I can provide detailed analytics for:

• **Sprint Analytics**: Completion rates, team performance, bug analysis
• **QA Metrics**: Quality scores, test results, bug trends  
• **Velocity Trends**: Sprint velocity history and trends
• **Team Performance**: Individual and team productivity metrics

Use the Analytics Dashboard buttons in the sidebar for detailed reports, or ask specific questions like:
- "Show me sprint analytics"
- "What's our QA quality score?"
- "How is our velocity trending?"
"""
                        else:
                            # Regular Jira query
                            jira_result = jira_tool(prompt, jira_client, board_id)
                        
                        if jira_result.get("items"):
                            response = f"Found {len(jira_result['items'])} Jira issues:\n\n"
                            for item in jira_result["items"]:
                                response += f"• **{item['key']}**: {item['summary']} ({item['status']}) - {item['assignee']}\n"
                                
                                # Add description if available and requested
                                if item.get('description') and ('description' in prompt.lower() or 'details' in prompt.lower() or 'full' in prompt.lower()):
                                    # Clean up the description (remove markdown formatting)
                                    description = item['description']
                                    description = description.replace('*', '').replace('\n', '\n  ')
                                    response += f"  **Description**:\n  {description}\n"
                                
                                response += f"  **Link**: {item['link']}\n\n"
                            response += f"JQL: `{jira_result['jql']}`"
                        else:
                            response = f"No Jira issues found. JQL: `{jira_result['jql']}`"
                    else:
                        response = "Jira is not properly configured. Please configure Jira settings in the sidebar."
                
                    if decision == "JIRA" and not st.session_state.get("jira_configured"):
                        # Jira query but Jira not configured - provide helpful response
                        if is_jira_related:
                            response = f"This appears to be a Jira-related question, but Jira is not configured. Please configure Jira settings in the sidebar to use Jira functionality.\n\n💡 **Pro Tip**: Once Jira is configured, I can help you with project management, sprint analytics, team performance metrics, and detailed issue tracking!"
                        else:
                            response = chat([{"role": "user", "content": prompt}])
                    
                    if decision == "GENERAL" or is_general_question:
                        # General chat - handle all other questions
                        if is_general_question:
                            response = f"{chat([{'role': 'user', 'content': prompt}])}\n\n💡 **Pro Tip**: I'm also designed to help with Jira project management and document analysis. Configure Jira in the sidebar to get insights about your projects, sprints, and team tasks. This way, I can help you with both general questions and work-related queries!"
                        else:
                            # General work-related chat
                            response = chat([{"role": "user", "content": prompt}])
                    
                    st.write(response)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                    
                    # Mark that we have responses to show export buttons
                    st.session_state.show_export_buttons = True
                
                except Exception as e:
                    error_msg = f"Error: {str(e)}"
                    st.error(error_msg)
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})

    # Export functionality
    if st.session_state.messages and not st.session_state.get("export_downloaded", False):
        st.markdown("---")
        st.subheader("📊 Export Options")
        
        col1, col2, col3 = st.columns([1, 1, 2])
        
        with col1:
            if st.button("📄 Export PDF", use_container_width=True):
                if st.session_state.messages:
                    try:
                        export_file = export_to_pdf(st.session_state.messages)
                        if export_file:
                            st.session_state.current_export_file = export_file
                            st.session_state.export_file_created = True
                            st.success("PDF exported successfully!")
                    except Exception as e:
                        st.error(f"PDF export failed: {e}")
        
        with col2:
            if st.button("📊 Export PPTX", use_container_width=True):
                if st.session_state.messages:
                    try:
                        export_file = export_to_powerpoint(st.session_state.messages)
                        if export_file:
                            st.session_state.current_export_file = export_file
                            st.session_state.export_file_created = True
                            st.success("PowerPoint exported successfully!")
                    except Exception as e:
                        st.error(f"PowerPoint export failed: {e}")
        
        with col3:
            if st.session_state.export_file_created and st.session_state.current_export_file:
                with open(st.session_state.current_export_file, "rb") as file:
                    file_extension = st.session_state.current_export_file.split('.')[-1].upper()
                    st.download_button(
                        label=f"⬇️ Download {file_extension}",
                        data=file.read(),
                        file_name=st.session_state.current_export_file,
                        mime="application/octet-stream",
                        use_container_width=True
                    )
                    if st.button("🗑️ Clear Export", use_container_width=True):
                        st.session_state.export_downloaded = True
                        st.session_state.export_file_created = False
                        st.session_state.current_export_file = None
                        st.rerun()
# Sidebar info - FIXED VERSION
with st.sidebar:
    # Clean company branding section
    st.markdown("""
    <div style="text-align: center; padding: 15px; margin-bottom: 25px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 12px; box-shadow: 0 4px 15px rgba(0,0,0,0.15);">
        <h3 style="color: white; margin: 0; font-size: 1.3rem; font-weight: bold; text-shadow: 1px 1px 2px rgba(0,0,0,0.3);">TAO DIGITAL SOLUTIONS</h3>
        <p style="color: rgba(255,255,255,0.9); margin: 5px 0 0 0; font-size: 1rem;">Integration Hub</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Integration Hub Section
    st.header("🔗 Integrations")
    
    # Integration Status Overview
    st.subheader("📊 Connection Status")
    
    # Jira Integration Card
    jira_status = "🟢 Connected" if st.session_state.get("jira_configured") else "🔴 Not Connected"
    st.markdown(f"""
    <div style="background: {'#d4edda' if st.session_state.get('jira_configured') else '#f8d7da'}; border: 1px solid {'#c3e6cb' if st.session_state.get('jira_configured') else '#f5c6cb'}; border-radius: 8px; padding: 12px; margin-bottom: 10px;">
        <div style="display: flex; align-items: center; justify-content: space-between;">
            <div>
                <strong>📋 Jira</strong><br>
                <small style="color: #6c757d;">Project Management</small>
            </div>
            <div style="font-size: 0.9rem;">{jira_status}</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Confluence Integration Card (Placeholder for future)
    st.markdown("""
    <div style="background: #e2e3e5; border: 1px solid #d6d8db; border-radius: 8px; padding: 12px; margin-bottom: 10px;">
        <div style="display: flex; align-items: center; justify-content: space-between;">
            <div>
                <strong>📚 Confluence</strong><br>
                <small style="color: #6c757d;">Knowledge Management</small>
            </div>
            <div style="font-size: 0.9rem; color: #6c757d;">🚧 Coming Soon</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # GitHub Integration Card (Placeholder for future)
    st.markdown("""
    <div style="background: #e2e3e5; border: 1px solid #d6d8db; border-radius: 8px; padding: 12px; margin-bottom: 20px;">
        <div style="display: flex; align-items: center; justify-content: space-between;">
            <div>
                <strong>🐙 GitHub</strong><br>
                <small style="color: #6c757d;">Code Repository</small>
            </div>
            <div style="font-size: 0.9rem; color: #6c757d;">🚧 Coming Soon</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Jira Configuration Section
    st.markdown("---")
    st.subheader("⚙️ Configure Jira")
    
    # Jira Configuration Button
    if st.button("🔧 Configure Jira", use_container_width=True, type="primary"):
        st.session_state.show_jira_config = True
    
    # Jira Configuration Modal/Popup
    if st.session_state.get("show_jira_config", False):
        st.markdown("---")
        st.markdown("### 🔧 Jira Configuration")
        
        # Add helpful instructions
        with st.expander("ℹ️ How to get Jira API Token", expanded=True):
            st.markdown("""
            **To configure Jira, you need an API Token:**
            
            1. Go to [Atlassian Account Settings](https://id.atlassian.com/manage-profile/security/api-tokens)
            2. Click "Create API token"
            3. Enter a label (e.g., "Integration Hub")
            4. Copy the generated token
            5. Use this token in the "API Token" field below
            
            **Note:** Use API token, not your login password!
            """)
        
        # Configuration form
        jira_url = st.text_input(
            "🌐 Jira URL", 
            placeholder="https://your-domain.atlassian.net",
            help="Your Jira instance URL"
        )
        jira_email = st.text_input(
            "📧 Email", 
            placeholder="your-email@example.com",
            help="Your Atlassian account email"
        )
        jira_token = st.text_input(
            "🔑 API Token", 
            type="password", 
            placeholder="your-api-token",
            help="API token from Atlassian account settings"
        )
        jira_board = st.text_input(
            "📋 Board ID (Optional)", 
            placeholder="123", 
            help="Find this in your Jira board URL"
        )
        
        # Configuration buttons
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            if st.button("✅ Save & Connect", use_container_width=True):
                if jira_url and jira_email and jira_token:
                    try:
                        # Set environment variables
                        os.environ["JIRA_BASE_URL"] = jira_url
                        os.environ["JIRA_EMAIL"] = jira_email
                        os.environ["JIRA_API_TOKEN"] = jira_token
                        os.environ["JIRA_BOARD_ID"] = jira_board if jira_board else ""
                        
                        # Create Jira client
                        jira_config = JiraConfig(
                            base_url=jira_url.rstrip("/"),
                            email=jira_email,
                            api_token=jira_token,
                            board_id=jira_board if jira_board else ""
                        )
                        jira_client = JiraClient(jira_config)
                        
                        # Test connection
                        sprint = jira_client.get_current_sprint()
                        
                        # Store in session state
                        st.session_state.jira_client = jira_client
                        st.session_state.jira_board_id = jira_board if jira_board else None
                        st.session_state.jira_configured = True
                        
                        st.success("✅ Jira configured successfully!")
                        if sprint:
                            st.info(f"🏃 Current sprint: {sprint.get('name', 'Unknown')}")
                        elif not jira_board:
                            st.info("ℹ️ Board ID not provided - sprint features disabled")
                        
                        # Close the configuration modal
                        st.session_state.show_jira_config = False
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"❌ Jira configuration failed: {e}")
                else:
                    st.warning("⚠️ Please fill in required fields: URL, Email, and API Token")
        
        with col2:
            if st.button("❌ Cancel", use_container_width=True):
                st.session_state.show_jira_config = False
                st.rerun()
        
        with col3:
            if st.button("🗑️ Clear Chat", use_container_width=True):
                st.session_state.messages = []
                st.session_state.document_text = ""
                st.session_state.last_uploaded_file = None
                st.session_state.clear_clicked = True
                st.session_state.show_export_buttons = False
                st.session_state.export_downloaded = False
                st.session_state.current_export_file = None
                st.session_state.export_file_created = False
                st.rerun()
    
    # Quick Actions Section
    st.markdown("---")
    st.subheader("⚡ Quick Actions")
    
    if st.session_state.get("jira_configured"):
        if st.button("📊 Sprint Analytics", use_container_width=True):
            try:
                jira_client = st.session_state.get("jira_client")
                board_id = st.session_state.get("jira_board_id")
                
                with st.spinner("Analyzing sprint data..."):
                    analytics = get_sprint_analytics(jira_client, board_id)
                
                if analytics:
                    st.success(f"📊 **{analytics['sprint_name']}** Analytics")
                    
                    # Key Metrics
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Issues", analytics['total_issues'])
                    with col2:
                        st.metric("Completed", analytics['completed_issues'])
                    with col3:
                        st.metric("Remaining", analytics['remaining_issues'])
                    
                    # Quality Score
                    quality_score = calculate_quality_score(analytics['bug_metrics'], analytics['test_metrics'])
                    st.metric("Quality Score", f"{quality_score}%")
                    
                    # Add to chat
                    analytics_summary = f"""
📊 **Sprint Analytics Summary**
- **Sprint**: {analytics['sprint_name']}
- **Total Issues**: {analytics['total_issues']}
- **Completed**: {analytics['completed_issues']}
- **Remaining**: {analytics['remaining_issues']}
- **Quality Score**: {quality_score}%

**Bug Analysis**: {analytics['bug_metrics']['total_bugs']} total bugs, {analytics['bug_metrics']['closed_bugs']} closed
**Test Coverage**: {analytics['test_metrics']['total_tests']} test cases, {analytics['test_metrics']['passed_tests']} passed
                    """
                    
                    st.session_state.messages.append({"role": "assistant", "content": analytics_summary})
                    st.rerun()
            except Exception as e:
                st.error(f"Analytics failed: {e}")
    else:
        st.info("🔗 Configure Jira to access analytics and project insights")
    
    # Check if we have previous connection details
    if "previous_jira_url" not in st.session_state:
        st.session_state.previous_jira_url = ""
    if "previous_jira_email" not in st.session_state:
        st.session_state.previous_jira_email = ""
    if "previous_jira_token" not in st.session_state:
        st.session_state.previous_jira_token = ""
    if "previous_jira_board" not in st.session_state:
        st.session_state.previous_jira_board = ""
    
    # Use Previous Connection button
    if (st.session_state.previous_jira_url and 
        st.session_state.previous_jira_email and 
        st.session_state.previous_jira_token):
        if st.button("🔄 Use Previous Connection", use_container_width=True, help="Use your last successful Jira connection"):
            try:
                # Create Jira client with previous values
                jira_config = JiraConfig(
                    base_url=st.session_state.previous_jira_url.rstrip("/"),
                    email=st.session_state.previous_jira_email,
                    api_token=st.session_state.previous_jira_token,
                    board_id=st.session_state.previous_jira_board if st.session_state.previous_jira_board else ""
                )
                jira_client = JiraClient(jira_config)
                
                # Test connection
                sprint = jira_client.get_current_sprint()
                
                # Store in session state
                st.session_state.jira_client = jira_client
                st.session_state.jira_board_id = st.session_state.previous_jira_board if st.session_state.previous_jira_board else None
                st.session_state.jira_configured = True
                
                st.success("✅ Previous Jira connection restored successfully!")
                if sprint:
                    st.info(f"🏃 Current sprint: {sprint.get('name', 'Unknown')}")
                elif not st.session_state.previous_jira_board:
                    st.info("ℹ️ Board ID not provided - sprint features disabled")
                    
            except Exception as e:
                error_msg = str(e)
                st.error(f"❌ Previous connection failed: {error_msg}")
                st.info("💡 Please configure Jira manually with fresh credentials")
    
    # Cleaner form layout with previous values as defaults
    jira_url = st.text_input(
        "🌐 Jira URL", 
        value=st.session_state.previous_jira_url,
        placeholder="https://your-domain.atlassian.net"
    )
    jira_email = st.text_input(
        "📧 Email", 
        value=st.session_state.previous_jira_email,
        placeholder="your-email@example.com"
    )
    jira_token = st.text_input(
        "🔑 API Token", 
        type="password", 
        value=st.session_state.previous_jira_token,
        placeholder="your-api-token"
    )
    jira_board = st.text_input(
        "📋 Board ID (Optional)", 
        value=st.session_state.previous_jira_board,
        placeholder="123", 
        help="Find this in your Jira board URL"
    )
    
    # Configuration button with better styling
    col1, col2 = st.columns([3, 1])
    with col1:
        if st.button("🔗 Configure Jira", use_container_width=True):
            if jira_url and jira_email and jira_token:
                try:
                    # Set environment variables
                    os.environ["JIRA_BASE_URL"] = jira_url
                    os.environ["JIRA_EMAIL"] = jira_email
                    os.environ["JIRA_API_TOKEN"] = jira_token
                    os.environ["JIRA_BOARD_ID"] = jira_board if jira_board else ""
                    
                    # Debug: Show what we're setting
                    st.info(f"🔧 Configuring Jira with URL: {jira_url}")
                    
                    # Create Jira client with explicit values
                    jira_config = JiraConfig(
                        base_url=jira_url.rstrip("/"),
                        email=jira_email,
                        api_token=jira_token,
                        board_id=jira_board if jira_board else ""
                    )
                    jira_client = JiraClient(jira_config)
                    
                    # Test connection
                    sprint = jira_client.get_current_sprint()
                    
                    # Store in session state
                    st.session_state.jira_client = jira_client
                    st.session_state.jira_board_id = jira_board if jira_board else None
                    st.session_state.jira_configured = True
                    
                    # Save connection details for future use
                    st.session_state.previous_jira_url = jira_url
                    st.session_state.previous_jira_email = jira_email
                    st.session_state.previous_jira_token = jira_token
                    st.session_state.previous_jira_board = jira_board if jira_board else ""
                    
                    st.success("✅ Jira configured successfully!")
                    if sprint:
                        st.info(f"🏃 Current sprint: {sprint.get('name', 'Unknown')}")
                    elif not jira_board:
                        st.info("ℹ️ Board ID not provided - sprint features disabled")
                        
                except Exception as e:
                    error_msg = str(e)
                    st.error(f"❌ Jira configuration failed: {error_msg}")
                    
                    # Show more specific error information
                    if "401" in error_msg:
                        st.error("🔐 **Authentication Error**: Check your API token and email")
                    elif "404" in error_msg:
                        st.error("🔍 **Not Found**: Check your Jira URL and Board ID")
                    elif "your-domain" in error_msg:
                        st.error("🌐 **URL Error**: Make sure you're using the correct Jira URL")
                    
                    # Show debug info
                    st.info(f"🔧 **Debug Info**:")
                    st.info(f"• URL: {jira_url}")
                    st.info(f"• Email: {jira_email}")
                    st.info(f"• Board ID: {jira_board}")
                    st.info(f"• Token: {'Set' if jira_token else 'Not set'}")
        else:
                st.warning("⚠️ Please fill in required fields: URL, Email, and API Token")
    
    with col2:
        # Clear Chat button - Always visible near Configure Jira
        if st.button("🗑️", help="Clear chat", key="clear_chat_sidebar"):
            st.session_state.messages = []
            st.session_state.document_text = ""
            st.session_state.last_uploaded_file = None
            st.session_state.clear_clicked = True
            st.session_state.show_export_buttons = False
            st.session_state.export_downloaded = False
            st.session_state.current_export_file = None
            st.session_state.export_file_created = False
            st.rerun()
    
    # Analytics Dashboard Section
    st.markdown("---")  # Separator line
    st.header("📊 Analytics Dashboard")
    
    if st.session_state.get("jira_configured"):
        # Sprint Analytics
        if st.button("📈 Sprint Analytics", use_container_width=True):
            try:
                jira_client = st.session_state.get("jira_client")
                board_id = st.session_state.get("jira_board_id")
                
                with st.spinner("Analyzing sprint data..."):
                    analytics = get_sprint_analytics(jira_client, board_id)
                    
                if analytics:
                    st.success(f"📊 **{analytics['sprint_name']}** Analytics")
                    
                    # Key Metrics
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Issues", analytics['total_issues'])
                    with col2:
                        st.metric("Completed", analytics['completed_issues'])
                    with col3:
                        st.metric("Completion Rate", f"{analytics['completion_rate']}%")
                    
                    # Progress Bar
                    progress = analytics['completion_rate'] / 100
                    st.progress(progress)
                    
                    # Team Performance
                    if analytics['team_performance']:
                        st.subheader("👥 Team Performance")
                        for member, perf in analytics['team_performance'].items():
                            completion_rate = (perf['completed'] / perf['total'] * 100) if perf['total'] > 0 else 0
                            st.write(f"**{member}**: {perf['completed']}/{perf['total']} completed ({completion_rate:.1f}%)")
                    
                    # Bug Analysis
                    if analytics['bug_analysis']['total_bugs'] > 0:
                        st.subheader("🐛 Bug Analysis")
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Total Bugs", analytics['bug_analysis']['total_bugs'])
                        with col2:
                            st.metric("Critical Bugs", analytics['bug_analysis']['critical_bugs'])
                        with col3:
                            st.metric("Resolved", analytics['bug_analysis']['resolved_bugs'])
                else:
                    st.warning("No sprint data available")
            except Exception as e:
                st.error(f"Error getting analytics: {e}")
        
        # QA Metrics
        if st.button("🔍 QA Metrics", use_container_width=True):
            try:
                jira_client = st.session_state.get("jira_client")
                board_id = st.session_state.get("jira_board_id")
                
                with st.spinner("Analyzing QA data..."):
                    qa_metrics = generate_qa_metrics(jira_client, board_id)
                
                if qa_metrics:
                    st.success("🔍 **QA Metrics**")
                    
                    # Quality Score
                    quality_score = qa_metrics['quality_score']
                    st.metric("Quality Score", f"{quality_score:.1f}/100")
                    
                    # Bug Metrics
                    bug_metrics = qa_metrics['bug_metrics']
                    st.subheader("🐛 Bug Metrics")
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Total Bugs", bug_metrics['total_bugs'])
                    with col2:
                        st.metric("Critical", bug_metrics['critical_bugs'])
                    with col3:
                        st.metric("High Priority", bug_metrics['high_priority_bugs'])
                    with col4:
                        st.metric("Resolved", bug_metrics['resolved_bugs'])
                    
                    # Test Metrics
                    test_metrics = qa_metrics['test_metrics']
                    if test_metrics['total_tests'] > 0:
                        st.subheader("🧪 Test Metrics")
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Total Tests", test_metrics['total_tests'])
                        with col2:
                            st.metric("Passed", test_metrics['passed_tests'])
                        with col3:
                            st.metric("Failed", test_metrics['failed_tests'])
                        
                        # Test Success Rate
                        success_rate = (test_metrics['passed_tests'] / test_metrics['total_tests'] * 100) if test_metrics['total_tests'] > 0 else 0
                        st.metric("Test Success Rate", f"{success_rate:.1f}%")
                else:
                    st.warning("No QA data available")
            except Exception as e:
                st.error(f"Error getting QA metrics: {e}")
        
        # Velocity Trends
        if st.button("📈 Velocity Trends", use_container_width=True):
            try:
                jira_client = st.session_state.get("jira_client")
                board_id = st.session_state.get("jira_board_id")
                
                with st.spinner("Analyzing velocity trends..."):
                    velocity_data = get_velocity_trends(jira_client, board_id)
                
                if velocity_data:
                    st.success("📈 **Velocity Trends**")
                    
                    # Create velocity chart
                    import pandas as pd
                    df = pd.DataFrame(velocity_data)
                    
                    # Display velocity metrics
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        avg_velocity = df['velocity'].mean()
                        st.metric("Avg Velocity", f"{avg_velocity:.1f}")
                    with col2:
                        latest_velocity = df['velocity'].iloc[-1]
                        st.metric("Latest Velocity", latest_velocity)
                    with col3:
                        velocity_trend = "📈" if df['velocity'].iloc[-1] > df['velocity'].iloc[-2] else "📉"
                        st.metric("Trend", velocity_trend)
                    
                    # Velocity table
                    st.subheader("📊 Sprint Velocity History")
                    st.dataframe(df, use_container_width=True)
                else:
                    st.warning("No velocity data available")
            except Exception as e:
                st.error(f"Error getting velocity trends: {e}")
    
    # Status indicator with cleaner design
    st.markdown("---")  # Separator line
    
    if st.session_state.get("jira_configured"):
        st.markdown("""
        <div style="padding: 10px; background: #d4edda; border-radius: 6px; border-left: 4px solid #28a745;">
            <strong style="color: #155724;">✅ Jira Connected</strong>
        </div>
        """, unsafe_allow_html=True)
        
        if st.button("🧪 Test Jira Connection", use_container_width=True):
            try:
                jira_client = st.session_state.get("jira_client")
                board_id = st.session_state.get("jira_board_id")
                
                # Test with a simple query
                result = jira_tool("show me recent issues", jira_client, board_id)
                if result.get("items"):
                    st.success(f"✅ Found {len(result['items'])} issues")
                else:
                    st.warning("⚠️ Connection working but no issues found")
                    
                # Test specific ticket search
                st.info("🔍 Testing CCM-283 search...")
                ticket_result = jira_tool("CCM-283", jira_client, board_id)
                if ticket_result and ticket_result.get("items"):
                    st.success(f"✅ CCM-283 found! Status: {ticket_result['items'][0]['status']}")
                else:
                    st.error(f"❌ CCM-283 not found: {ticket_result.get('message', 'Unknown error')}")
                    
            except Exception as e:
                st.error(f"❌ Test failed: {e}")
    else:
        st.markdown("""
        <div style="padding: 10px; background: #fff3cd; border-radius: 6px; border-left: 4px solid #ffc107;">
            <strong style="color: #856404;">⚠️ Jira Not Configured</strong>
            <br><small>Configure above to enable Jira features</small>
        </div>
        """, unsafe_allow_html=True)
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; padding: 10px; color: #6c757d; font-size: 0.8rem;">
        <em>Powered by TAO Digital Solutions</em>
    </div>
    """, unsafe_allow_html=True)
