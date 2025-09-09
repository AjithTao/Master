from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import tempfile
import json
import io
import re
from datetime import datetime
import logging
import sys
import asyncio
from contextlib import asynccontextmanager

# Regex for detecting Jira ticket keys (case-insensitive)
ISSUE_KEY_RE = r"\b([A-Z][A-Z0-9]+-\d+)\b"

# Optional imports with fallback
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Ensure backend and src are importable (backend first so `utils` resolves to backend package)
BASE_DIR = os.path.dirname(__file__)
SRC_DIR = os.path.abspath(os.path.join(BASE_DIR, '..', 'src'))
if BASE_DIR not in sys.path:
    sys.path.insert(0, BASE_DIR)
if SRC_DIR not in sys.path:
    # insert after backend so backend utils wins over src/utils.py
    sys.path.insert(1, SRC_DIR)
from jira_client import JiraClient
from auth import JiraConfig
from ai_engine import AdvancedAIEngine, AIInsight
from query_processor import AdvancedQueryProcessor
from analytics_engine import AdvancedAnalyticsEngine
from enhanced_jql_processor import EnhancedJQLProcessor, ResponseFormat
from advanced_chatbot import AdvancedChatbotEngine, QueryIntent
from utils.jql_training_loader import JQLTrainingLoader
from utils.enhanced_jql_training_loader import EnhancedJQLTrainingLoader
from confluence_client import ConfluenceClient, ConfluenceConfig
from utils.metrics_utils import (
    parse_sprint_number, to_dt, find_story_points_field, sum_story_points,
    is_done, forecast_velocity, parse_sprint_metrics, forecast_velocity_wrapper,
    analyze_bandwidth, generate_insights
)
from datetime import datetime, timezone
from intent_router import handle_user_query as ai_handle
from ai_summarizer import summarize as ai_summarize

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Application state
class AppState:
    def __init__(self):
        self.jira_configured = False
        self.jira_client = None
        self.jira_config = None
        self.jira_board_id = None
        # Confluence
        self.confluence_configured = False
        self.confluence_client = None
        self.confluence_config = None
        # Cached Jira projects for the current session
        self.jira_projects = {}  # key -> name
        self.jira_project_keys = []  # list of keys for quick access
        self.messages = []
        self.export_files = {}
        self.ai_engine = None
        self.query_processor = None
        self.analytics_engine = None
        self.enhanced_jql_processor = None
        self.advanced_chatbot = None
        self.max_messages = 1000  # Keep last 1000 messages
        # Listing context (for pagination/export continuity)
        self.last_list_jql: str | None = None
        self.last_list_offset: int = 0
        self.last_list_total: int = 0
        self.last_list_keys_shown: list[str] = []

app_state = AppState()

# Initialize enhanced training loader with slot-based support
TRAINING_PATH = os.path.join(os.path.dirname(__file__), "data", "jira_ai_training_pack.json")
jql_loader = EnhancedJQLTrainingLoader(TRAINING_PATH)

# ------------------------------
# Friendly response helpers
# ------------------------------
def build_issue_url(key: str) -> str:
    base = app_state.jira_config.base_url if app_state.jira_config else ""
    return f"{base}/browse/{key}" if base else key

# (moved) Jira issue details endpoint is defined below after app initialization

async def fetch_top_items(jql: str, limit: int = 5) -> List[Dict[str, str]]:
    items: List[Dict[str, str]] = []
    try:
        # Ensure client initialized
        if not app_state.jira_client or not app_state.jira_client._client:
            await app_state.jira_client.initialize()
        result = await app_state.jira_client.search(jql, max_results=limit, fields=["key","summary"])
        for issue in result.get("issues", []):
            key = issue.get("key", "")
            summary = (issue.get("fields", {}) or {}).get("summary", "")
            items.append({
                "key": key,
                "summary": summary,
                "url": build_issue_url(key)
            })
    except Exception as e:
        logger.warning(f"fetch_top_items failed for JQL '{jql}': {e}")
    return items

def format_friendly_response(kind: str, count: int, context: Dict[str, Any]) -> str:
    """Create concise, friendly responses with light emojis.
    Supported kinds: closed_bugs, project_closed_stories, project_open_stories,
    assignee_open_stories, assignee_open_tasks, generic
    """
    project = context.get("project")
    assignee = context.get("assignee")
    # Map kinds to message templates
    if kind == "closed_bugs":
        return f"✅ I found {count} closed bugs in total. These are all completed Bug issues."
    if kind == "project_closed_stories" and project:
        return f"🎉 In the {project} project, {count} stories have been completed. Great progress!"
    if kind == "project_open_stories" and project:
        return f"🔄 The {project} project currently has {count} open stories. These are ongoing or pending work items."
    if kind == "assignee_open_stories" and assignee:
        return f"👤 {assignee} has {count} open stories assigned. Still in progress or waiting to be completed."
    if kind == "assignee_open_tasks" and assignee:
        return f"👤 {assignee} has {count} open tasks assigned. These need attention."
    # Fallback generic
    return f"📊 Found {count} items for your query."

def create_error_response(error_type: str, details: str = "", status_code: int = 500) -> Dict[str, Any]:
    """Create consistent error responses"""
    return {
        "success": False,
        "error": error_type,
        "details": details,
        "status_code": status_code
    }

def create_success_response(data: Any = None, message: str = "Success") -> Dict[str, Any]:
    """Create consistent success responses"""
    response = {
        "success": True,
        "message": message
    }
    if data is not None:
        response["data"] = data
    return response

def mask_email(email: str) -> str:
    """Mask email address for security"""
    if not email or "@" not in email:
        return email
    
    local, domain = email.split("@", 1)
    if len(local) <= 2:
        masked_local = "*" * len(local)
    else:
        masked_local = local[:2] + "*" * (len(local) - 2)
    
    return f"{masked_local}@{domain}"

def manage_message_history(app_state: AppState, message: Dict[str, Any]) -> None:
    """Manage message history to prevent memory bloat"""
    app_state.messages.append(message)
    
    # Keep only the last max_messages
    if len(app_state.messages) > app_state.max_messages:
        app_state.messages = app_state.messages[-app_state.max_messages:]

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan manager"""
    logger.info("🚀 Starting Leadership Management Tool API")
    yield
    logger.info("🛑 Shutting down Leadership Management Tool API")

app = FastAPI(
    title="Leadership Management Tool API", 
    version="1.0.0", 
    description="AI-powered leadership analytics and project management insights",
    lifespan=lifespan,
    openapi_tags=[
        {
            "name": "Chat",
            "description": "AI chat and conversation endpoints"
        },
        {
            "name": "JIRA",
            "description": "JIRA integration and analytics endpoints"
        },
        {
            "name": "Export",
            "description": "Data export and reporting endpoints"
        },
        {
            "name": "Analytics",
            "description": "Advanced analytics and insights endpoints"
        }
    ]
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/api/jira/issue/{issue_key}", tags=["JIRA"], summary="Get Jira issue details")
async def get_issue_details(issue_key: str):
    """Return detailed Jira issue fields for UI consumption."""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")

    try:
        issue = await app_state.jira_client.get_issue(issue_key, fields=[
            "summary", "status", "issuetype", "priority", "assignee", "reporter", "project", "created", "updated"
        ])
        if not issue:
            raise HTTPException(status_code=404, detail="Issue not found")

        f = issue.get("fields", {})
        data = {
            "key": issue.get("key", issue_key),
            "summary": f.get("summary"),
            "status": (f.get("status") or {}).get("name"),
            "type": (f.get("issuetype") or {}).get("name"),
            "priority": (f.get("priority") or {}).get("name"),
            "assignee": {
                "name": (f.get("assignee") or {}).get("displayName") if f.get("assignee") else None,
                "email": (f.get("assignee") or {}).get("emailAddress") if f.get("assignee") else None,
            },
            "reporter": {
                "name": (f.get("reporter") or {}).get("displayName") if f.get("reporter") else None,
                "email": (f.get("reporter") or {}).get("emailAddress") if f.get("reporter") else None,
            },
            "project": (f.get("project") or {}).get("key"),
            "created": f.get("created"),
            "updated": f.get("updated"),
            "url": build_issue_url(issue.get("key", issue_key))
        }

        return {"success": True, "issue": data}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/export/excel", tags=["JIRA"], summary="Export Jira items to Excel")
async def export_jira_items_to_excel(request: Dict[str, Any]):
    """Export Jira items to Excel format"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Prefer last_list_jql for user-triggered exports if explicit JQL not provided
        jql = request.get("jql") or getattr(app_state, 'last_list_jql', None) or "project is not EMPTY"
        max_results = request.get("max_results", 1000)
        
        # Get items from Jira
        result = await app_state.jira_client.search(jql, max_results=max_results, fields=[
            "key", "summary", "status", "issuetype", "priority", "assignee", "reporter", "created", "updated", "project"
        ])
        issues = result.get("issues", [])
        
        if not issues:
            raise HTTPException(status_code=404, detail="No items found to export")
        
        # Create Excel file
        import pandas as pd
        from io import BytesIO
        import base64
        
        # Prepare data for Excel
        data = []
        for issue in issues:
            fields = issue.get("fields", {})
            assignee = fields.get("assignee", {})
            reporter = fields.get("reporter", {})
            
            data.append({
                "Ticket": issue.get("key", ""),
                "Summary": fields.get("summary", ""),
                "Status": fields.get("status", {}).get("name", ""),
                "Type": fields.get("issuetype", {}).get("name", ""),
                "Priority": fields.get("priority", {}).get("name", ""),
                "Assignee": assignee.get("displayName", "Unassigned") if assignee else "Unassigned",
                "Reporter": reporter.get("displayName", "Unknown") if reporter else "Unknown",
                "Project": fields.get("project", {}).get("key", ""),
                "Created": fields.get("created", ""),
                "Updated": fields.get("updated", "")
            })
        
        # Create DataFrame and export to Excel
        df = pd.DataFrame(data)
        excel_buffer = BytesIO()
        df.to_excel(excel_buffer, index=False, sheet_name="Jira Items")
        excel_buffer.seek(0)
        
        # Encode to base64 for download
        excel_data = base64.b64encode(excel_buffer.getvalue()).decode()
        filename = f"jira_items_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        return {
            "success": True,
            "filename": filename,
            "data": excel_data,
            "count": len(issues),
            "message": f"Excel file created with {len(issues)} items"
        }
        
    except Exception as e:
        logger.error(f"Error exporting to Excel: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Pydantic models
class JiraConfigRequest(BaseModel):
    url: str
    email: str
    api_token: str
    board_id: Optional[str] = None

class ConfluenceConfigRequest(BaseModel):
    url: str
    email: str
    api_token: str

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    messages: List[ChatMessage] = []

class AssigneeCountRequest(BaseModel):
    assignee: str
    issuetype: Optional[str] = None
    project: Optional[str] = None
    status_category: Optional[str] = None  # e.g., "Done", "To Do", "In Progress"
    not_done: Optional[bool] = True  # default to open items (!= Done)

# ---------------------
# Metrics Endpoints
# ---------------------

class MetricsRequest(BaseModel):
    project: Optional[str] = None
    window_days: int = 14


def _project_clause(project: Optional[str]) -> str:
    return f'project = "{project}" AND ' if project else ''


@app.post("/api/metrics/completed", tags=["Metrics"], summary="Completed stories and bugs")
async def metrics_completed(req: MetricsRequest):
    if not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    wnd = max(7, min(req.window_days or 14, 45))
    jql_now = f"{_project_clause(req.project)}issuetype in (Story, Bug) AND statusCategory = Done AND updated >= -{wnd}d"
    jql_prev = f"{_project_clause(req.project)}issuetype in (Story, Bug) AND statusCategory = Done AND updated >= -{2*wnd}d AND updated < -{wnd}d"
    now = await app_state.jira_client.search(jql_now, max_results=0)
    prev = await app_state.jira_client.search(jql_prev, max_results=0)
    now_count = int(now.get('total', 0))
    prev_count = int(prev.get('total', 0))
    delta = now_count - prev_count
    pct = (delta / prev_count * 100.0) if prev_count > 0 else (100.0 if now_count > 0 else 0.0)
    return {"count": now_count, "trend": {"delta": delta, "percent": round(pct, 1)}}


@app.post("/api/metrics/blockers", tags=["Metrics"], summary="Unresolved blockers by priority")
async def metrics_blockers(req: MetricsRequest):
    if not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    jql = f"{_project_clause(req.project)}status = Blocked AND resolution is EMPTY"
    result = await app_state.jira_client.search(jql, max_results=100, fields=["priority"]) 
    issues = result.get("issues", [])
    by_priority: Dict[str, int] = {}
    for it in issues:
        name = ((it.get("fields", {}).get("priority") or {}).get("name")) or "Unknown"
        by_priority[name] = by_priority.get(name, 0) + 1
    total = sum(by_priority.values())
    return {"total": total, "by_priority": by_priority}


@app.post("/api/metrics/contributors", tags=["Metrics"], summary="Top contributors in window")
async def metrics_contributors(req: MetricsRequest):
    if not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    wnd = max(7, min(req.window_days or 14, 90))
    jql = f"{_project_clause(req.project)}statusCategory = Done AND updated >= -{wnd}d"
    result = await app_state.jira_client.search(jql, max_results=200, fields=["assignee"]) 
    issues = result.get("issues", [])
    by_user: Dict[str, int] = {}
    for it in issues:
        assignee = (it.get("fields", {}).get("assignee") or {}).get("displayName") or "Unassigned"
        by_user[assignee] = by_user.get(assignee, 0) + 1
    # Top 10
    top = sorted(by_user.items(), key=lambda x: x[1], reverse=True)[:10]
    return {"window_days": wnd, "leaders": [{"name": n, "closed": c} for n, c in top]}


@app.post("/api/metrics/resolution", tags=["Metrics"], summary="Average resolution time (days)")
async def metrics_resolution(req: MetricsRequest):
    if not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    wnd = max(14, min(req.window_days or 30, 120))
    # Use updated as a proxy for resolution timestamp for Done items
    jql = f"{_project_clause(req.project)}statusCategory = Done AND updated >= -{wnd}d"
    result = await app_state.jira_client.search(jql, max_results=200, fields=["created", "updated"]) 
    issues = result.get("issues", [])
    import datetime as _dt
    durations: List[float] = []
    for it in issues:
        f = it.get("fields", {})
        try:
            created = _dt.datetime.fromisoformat((f.get("created") or "").replace("Z", "+00:00"))
            updated = _dt.datetime.fromisoformat((f.get("updated") or "").replace("Z", "+00:00"))
            durations.append((updated - created).total_seconds() / 86400.0)
        except Exception:
            continue
    avg_days = round(sum(durations) / len(durations), 2) if durations else 0.0
    return {"avg_days": avg_days, "sample": len(durations)}


@app.post("/api/metrics/velocity", tags=["Metrics"], summary="Velocity comparison current vs prior window")
async def metrics_velocity(req: MetricsRequest):
    if not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    wnd = max(7, min(req.window_days or 14, 45))
    cur_jql = f"{_project_clause(req.project)}statusCategory = Done AND updated >= -{wnd}d"
    prev_jql = f"{_project_clause(req.project)}statusCategory = Done AND updated >= -{2*wnd}d AND updated < -{wnd}d"
    cur = await app_state.jira_client.search(cur_jql, max_results=0)
    prev = await app_state.jira_client.search(prev_jql, max_results=0)
    cur_count = int(cur.get('total', 0))
    prev_count = int(prev.get('total', 0))
    delta = cur_count - prev_count
    pct = (delta / prev_count * 100.0) if prev_count > 0 else (100.0 if cur_count > 0 else 0.0)
    return {"current": cur_count, "previous": prev_count, "change_percent": round(pct, 1)}


@app.post("/api/metrics/summary", tags=["Metrics"], summary="AI leadership insight from metrics bundle")
async def metrics_summary(payload: Dict[str, Any]):
    # payload example: { project, completed, blockers, contributors, resolution, velocity }
    project = payload.get("project") or "this project"
    completed = payload.get("completed", {}).get("count", 0)
    blockers = payload.get("blockers", {}).get("total", 0)
    top = payload.get("contributors", {}).get("leaders", [])
    top_str = ", ".join([f"{x['name']} ({x['closed']})" for x in top[:3]]) if top else ""
    avg_res = payload.get("resolution", {}).get("avg_days", 0)
    vel = payload.get("velocity", {})
    cur_v, prev_v, pct_v = vel.get("current", 0), vel.get("previous", 0), vel.get("change_percent", 0)

    # Deterministic base narrative
    base = (
        f"Velocity {('up' if pct_v>0 else 'down' if pct_v<0 else 'flat')} {abs(pct_v)}% (" 
        f"{cur_v} vs {prev_v}). Blockers: {blockers}. Completed: {completed}. "
        f"Avg resolution: {avg_res} days. Top: {top_str}."
    ).strip()

    # Optional LLM overlay via existing summarizer
    try:
        overlay = ai_summarize(
            user_query=f"Leadership insight for {project}",
            results=[{"text": base}],
            response_type="list",
        )
        return {"insight": overlay or base}
    except Exception:
        return {"insight": base}

# Helper functions
async def handle_jira_question(message: str, jira_client: JiraClient) -> str:
    """Handle Jira-related questions with intelligent parsing"""
    try:
        message_lower = message.lower()
        
        # Extract ticket key (e.g., CCM-283)
        ticket_match = re.search(r'([A-Z][A-Z0-9]+-\d+)', message, re.IGNORECASE)
        
        if ticket_match:
            ticket_key = ticket_match.group(1).upper()
            return await get_ticket_details(ticket_key, jira_client)
        
        # Intelligent question routing with comprehensive pattern matching
        
        # Handle assignee-related questions (enhanced patterns)
        assignee_patterns = [
            'worked', 'assigned', 'assignee', 'who is', 'who has', 'who did', 'who does',
            'owner', 'responsible', 'developer', 'engineer', 'programmer', 'tester',
            'ashwin', 'thyagarajan', 'john', 'doe', 'jane', 'smith', 'mike', 'johnson',
            'sarah', 'wilson', 'david', 'brown', 'lisa', 'davis', 'robert', 'miller',
            'emily', 'jones', 'chris', 'anderson', 'amanda', 'taylor'
        ]
        if any(word in message_lower for word in assignee_patterns):
            return await search_by_assignee(message, jira_client)
        
        # Handle project-related questions (enhanced patterns)
        project_patterns = [
            'project', 'projects', 'proj', 'initiative', 'program', 'portfolio', 'product',
            'application', 'app', 'system', 'module', 'component', 'service',
            'how many projects', 'project count', 'project summary', 'project status',
            'project overview', 'project breakdown', 'project details'
        ]
        if any(word in message_lower for word in project_patterns):
            return await get_project_info(jira_client)
        
        # Handle sprint-related questions (enhanced patterns)
        sprint_patterns = [
            'sprint', 'sprints', 'iteration', 'iterations', 'cycle', 'cycles',
            'current sprint', 'active sprint', 'running sprint', 'ongoing sprint',
            'live sprint', 'open sprint', 'closed sprint', 'completed sprint',
            'finished sprint', 'done sprint', 'ended sprint', 'started sprint',
            'sprint status', 'sprint info', 'sprint details', 'sprint overview',
            'sprint progress', 'sprint velocity', 'sprint burndown', 'sprint burnup'
        ]
        if any(word in message_lower for word in sprint_patterns):
            return await get_sprint_info(jira_client)
        
        # Handle status-related questions (enhanced patterns)
        status_patterns = [
            'status', 'state', 'stage', 'phase', 'step', 'progress', 'progression',
            'todo', 'to do', 'open', 'new', 'assigned', 'in progress', 'in-progress',
            'under review', 'review', 'testing', 'test', 'qa', 'ready for test',
            'ready for qa', 'ready for review', 'ready for deploy', 'deployed',
            'released', 'closed', 'resolved', 'fixed', 'completed', 'done', 'finished',
            'cancelled', 'blocked', 'on hold', 'pending', 'waiting', 'stalled', 'stuck',
            'delayed', 'status breakdown', 'status overview', 'status summary',
            'status count', 'status report', 'status analytics'
        ]
        if any(word in message_lower for word in status_patterns):
            return await get_status_info(message, jira_client)
        
        # Handle issue type questions (enhanced patterns)
        issue_type_patterns = [
            'story', 'stories', 'bug', 'bugs', 'defect', 'defects', 'task', 'tasks',
            'epic', 'epics', 'subtask', 'subtasks', 'improvement', 'improvements',
            'feature', 'features', 'requirement', 'requirements', 'user story', 'user stories',
            'issue type', 'issue types', 'type breakdown', 'type overview', 'type summary',
            'type count', 'type report', 'type analytics', 'story count', 'bug count',
            'task count', 'defect count', 'feature count'
        ]
        if any(word in message_lower for word in issue_type_patterns):
            return await get_issue_type_info(message, jira_client)
        
        # Handle general analytics (enhanced patterns)
        analytics_patterns = [
            'summary', 'summaries', 'overview', 'overviews', 'analytics', 'analysis',
            'dashboard', 'metrics', 'kpi', 'kpis', 'statistics', 'stats', 'data',
            'insights', 'breakdown', 'breakdowns', 'count', 'counts', 'counting',
            'total', 'totals', 'number', 'numbers', 'how many', 'how much',
            'quantity', 'quantities', 'amount', 'amounts', 'percentage', 'percent',
            'ratio', 'ratios', 'rate', 'rates', 'trend', 'trends', 'trending',
            'pattern', 'patterns', 'distribution', 'report', 'reports', 'reporting'
        ]
        if any(word in message_lower for word in analytics_patterns):
            return await get_general_analytics(jira_client)
        
        # Default response with suggestions
        return """Hey! I can help you find stuff in Jira. Here's what I can do:

🎫 **Find tickets:** "Tell me about CCM-283" or "What's up with CCM-123?"

👤 **Check who's working on what:** "What's Ashwin working on?" or "Who's got CCM-283?"

📁 **Project stuff:** "How many projects do we have?" or "What's our project status?"

🏃 **Sprint info:** "What's our current sprint?" or "Sprint details please"

📊 **Status check:** "What's in progress?" or "Show me completed work"

Just ask me naturally - I'll figure out what you need! 😊"""
    
    except Exception as e:
        logger.error(f"Jira question handling error: {e}")
        return f"Oops! I ran into an issue while processing your question: {str(e)}"

async def get_ticket_details(ticket_key: str, jira_client: JiraClient) -> str:
    """Get detailed information about a specific ticket"""
    try:
        jql = f"key = {ticket_key}"
        result = await jira_client.search(jql, max_results=1)
        
        if result.get('total', 0) > 0:
            issue = result['issues'][0]
            fields = issue.get('fields', {})
            
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            assignee_email = assignee.get('emailAddress', '') if assignee else ''
            
            status = fields.get('status', {}).get('name', 'Unknown')
            summary = fields.get('summary', 'No summary')
            issue_type = fields.get('issuetype', {}).get('name', 'Unknown')
            priority = fields.get('priority', {}).get('name', 'Unknown')
            project = fields.get('project', {}).get('name', 'Unknown')
            
            # Get additional fields if available
            description = fields.get('description', 'No description')
            created = fields.get('created', 'Unknown')
            updated = fields.get('updated', 'Unknown')
            
            # Create grounded summary with plain formatting as requested
            canonical_key = issue.get("key", ticket_key).upper()
            response_text = (
                "Hello !! Here are your details\n\n"
                f"{canonical_key}: {summary}\n\n"
                f"Status: {status}\n"
                f"Type: {issue_type}\n"
                f"Priority: {priority}\n"
                f"Assignee: {assignee_name}\n"
                f"Project: {project}"
            )
            
            if description and description != 'No description':
                # Truncate long descriptions
                desc_preview = description[:200] + "..." if len(description) > 200 else description
                response_text += f"\n📝 **Description:** {desc_preview}"
            
            return {
                "success": True,
                "response": response_text
            }
        else:
            return f"Oops! I couldn't find ticket {ticket_key}. Double-check the ticket number and try again! 🤔"
    
    except Exception as e:
        logger.error(f"Error getting ticket details: {e}")
        return f"Oops! I had trouble getting details for {ticket_key}: {str(e)}"

async def search_by_assignee(message: str, jira_client: JiraClient) -> str:
    """Search by assignee or project using advanced JQL generation with fuzzy matching"""
    try:
        from utils.advanced_jql_generator import AdvancedJQLGenerator
        
        generator = AdvancedJQLGenerator(jira_client)
        result = await generator.generate_jql(message)
        
        if result['success']:
            return result['response']
        else:
            return result['response']
            
    except Exception as e:
        logger.error(f"Error in search_by_assignee: {e}")
        return f"Sorry, I encountered an error while searching: {str(e)}"
async def get_project_info(jira_client: JiraClient) -> str:
    """Get project information and statistics"""
    try:
        # Get all projects
        result = await jira_client.search("project is not EMPTY", max_results=1000)
        issues = result.get('issues', [])
        
        projects = {}
        for issue in issues:
            project_key = issue.get('fields', {}).get('project', {}).get('key', 'Unknown')
            project_name = issue.get('fields', {}).get('project', {}).get('name', 'Unknown')
            
            if project_key not in projects:
                projects[project_key] = {
                    'name': project_name,
                    'stories': 0,
                    'bugs': 0,
                    'tasks': 0,
                    'total': 0
                }
            
            projects[project_key]['total'] += 1
            
            issue_type = issue.get('fields', {}).get('issuetype', {}).get('name', '').lower()
            if 'story' in issue_type:
                projects[project_key]['stories'] += 1
            elif 'bug' in issue_type or 'defect' in issue_type:
                projects[project_key]['bugs'] += 1
            elif 'task' in issue_type:
                projects[project_key]['tasks'] += 1
        
        response = f"Here's what I found across your **{len(projects)} projects**:\n\n"
        
        for project_key, data in projects.items():
            response += f"🏢 **{project_key}** - {data['name']}\n"
            response += f"   📊 Total: {data['total']} | 📖 Stories: {data['stories']} | 🐛 Bugs: {data['bugs']} | ✅ Tasks: {data['tasks']}\n\n"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting project info: {e}")
        return f"Sorry, I couldn't retrieve project information: {str(e)}"

async def get_sprint_info(jira_client: JiraClient) -> str:
    """Get current sprint information"""
    try:
        current_sprint = jira_client.get_current_sprint()
        if current_sprint:
            return f"""Here's your current sprint info:

🏃 **Sprint:** {current_sprint.get('name', 'Unknown')}
📊 **Status:** {current_sprint.get('state', 'Unknown')}
📅 **Started:** {current_sprint.get('startDate', 'Not set')}
📅 **Ends:** {current_sprint.get('endDate', 'Not set')}
🎯 **Goal:** {current_sprint.get('goal', 'No goal set')}"""
        else:
            return "Hmm, looks like there's no active sprint right now. Maybe check your sprint configuration? 🤔"
    
    except Exception as e:
        logger.error(f"Error getting sprint info: {e}")
        return f"Sorry, I couldn't retrieve sprint information: {str(e)}"

async def get_status_info(message: str, jira_client: JiraClient) -> str:
    """Get information about ticket statuses"""
    try:
        message_lower = message.lower()
        
        if 'todo' in message_lower or 'to do' in message_lower:
            jql = 'status = "To Do"'
            status_name = "To Do"
        elif 'in progress' in message_lower or 'progress' in message_lower:
            jql = 'status = "In Progress"'
            status_name = "In Progress"
        elif 'done' in message_lower or 'completed' in message_lower:
            jql = 'status = "Done"'
            status_name = "Done"
        else:
            # Get all statuses
            result = await jira_client.search("project is not EMPTY", max_results=100)
            issues = result.get('issues', [])
            
            status_counts = {}
            for issue in issues:
                status = issue.get('fields', {}).get('status', {}).get('name', 'Unknown')
                status_counts[status] = status_counts.get(status, 0) + 1
            
            response = "Here's your ticket status breakdown:\n\n"
            for status, count in sorted(status_counts.items()):
                response += f"🔸 **{status}:** {count} tickets\n"
            
            return response
        
        result = await jira_client.search(jql, max_results=50)
        issues = result.get('issues', [])
        
        response = f"Here are the tickets in **{status_name}** status ({len(issues)} total):\n\n"
        
        for issue in issues[:10]:  # Show max 10
            fields = issue.get('fields', {})
            key = issue.get('key', 'Unknown')
            summary = fields.get('summary', 'No summary')
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            
            response += f"🎫 **{key}** - {summary}\n"
            response += f"   👤 Assigned to: {assignee_name}\n\n"
        
        if len(issues) > 10:
            response += f"... and {len(issues) - 10} more tickets! 😊"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting status info: {e}")
        return f"Sorry, I couldn't retrieve status information: {str(e)}"

async def get_issue_type_info(message: str, jira_client: JiraClient) -> str:
    """Get information about issue types"""
    try:
        message_lower = message.lower()
        
        if 'story' in message_lower or 'stories' in message_lower:
            jql = 'issuetype = Story'
            issue_type = "Story"
        elif 'bug' in message_lower or 'bugs' in message_lower or 'defect' in message_lower:
            jql = 'issuetype in (Bug, Defect)'
            issue_type = "Bug/Defect"
        elif 'task' in message_lower or 'tasks' in message_lower:
            jql = 'issuetype = Task'
            issue_type = "Task"
        else:
            # Get all issue types
            result = await jira_client.search("project is not EMPTY", max_results=100)
            issues = result.get('issues', [])
            
            type_counts = {}
            for issue in issues:
                issue_type = issue.get('fields', {}).get('issuetype', {}).get('name', 'Unknown')
                type_counts[issue_type] = type_counts.get(issue_type, 0) + 1
            
            response = "Here's your issue type breakdown:\n\n"
            for issue_type, count in sorted(type_counts.items()):
                response += f"🔸 **{issue_type}:** {count} tickets\n"
            
            return response
        
        result = await jira_client.search(jql, max_results=50)
        issues = result.get('issues', [])
        
        response = f"Here are your **{issue_type}** tickets ({len(issues)} total):\n\n"
        
        for issue in issues[:10]:  # Show max 10
            fields = issue.get('fields', {})
            key = issue.get('key', 'Unknown')
            summary = fields.get('summary', 'No summary')
            status = fields.get('status', {}).get('name', 'Unknown')
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            
            response += f"🎫 **{key}** - {summary}\n"
            response += f"   📊 Status: {status} | 👤 Assigned to: {assignee_name}\n\n"
        
        if len(issues) > 10:
            response += f"... and {len(issues) - 10} more tickets! 😊"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting issue type info: {e}")
        return f"Sorry, I couldn't retrieve issue type information: {str(e)}"

async def get_general_analytics(jira_client: JiraClient) -> str:
    """Get general analytics and summary"""
    try:
        result = await jira_client.search("project is not EMPTY", max_results=1000)
        issues = result.get('issues', [])
        
        # Calculate statistics
        total_issues = len(issues)
        projects = set()
        assignees = set()
        status_counts = {}
        type_counts = {}
        
        for issue in issues:
            fields = issue.get('fields', {})
            
            # Count projects
            project_key = fields.get('project', {}).get('key', 'Unknown')
            projects.add(project_key)
            
            # Count assignees
            assignee = fields.get('assignee')
            if assignee and assignee.get('displayName'):
                assignees.add(assignee.get('displayName'))
            
            # Count statuses
            status = fields.get('status', {}).get('name', 'Unknown')
            status_counts[status] = status_counts.get(status, 0) + 1
            
            # Count types
            issue_type = fields.get('issuetype', {}).get('name', 'Unknown')
            type_counts[issue_type] = type_counts.get(issue_type, 0) + 1
        
        response = f"""Here's your Jira overview:

🎯 **Quick Stats:**
• Total Issues: {total_issues}
• Projects: {len(projects)}
• Team Members: {len(assignees)}

📊 **Status Breakdown:**
"""
        for status, count in sorted(status_counts.items()):
            response += f"• {status}: {count}\n"
        
        response += f"\n🏷️ **Issue Types:**\n"
        for issue_type, count in sorted(type_counts.items()):
            response += f"• {issue_type}: {count}\n"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting general analytics: {e}")
        return f"Sorry, I couldn't retrieve analytics information: {str(e)}"

async def handle_general_question(message: str) -> str:
    """Handle general questions with leadership focus"""
    if 'hello' in message or 'hi' in message:
        return """Hello! I'm your Leadership Quality Assistant for TAO Digital. I specialize in transforming project management data into actionable business insights.

I can help you with:
• Team performance analysis and capacity planning
• Project health assessment and risk identification  
• Sprint velocity trends and quality metrics
• Resource allocation and process optimization
• Strategic recommendations based on your data

What leadership insights would you like to explore today?"""
    
    elif 'help' in message:
        return """As your Leadership Quality Assistant, I provide strategic insights from your project data:

**TEAM PERFORMANCE:**
• "Show me Ajith's current workload and capacity"
• "Analyze team velocity trends for the last 3 sprints"
• "Who's overloaded and needs resource reallocation?"

**PROJECT HEALTH:**
• "What's the status of CCM project and any blockers?"
• "Identify risks in our current sprint"
• "Show me project progress against goals"

**QUALITY & PROCESS:**
• "Analyze our defect rates and quality trends"
• "What process improvements do you recommend?"
• "How can we optimize our sprint planning?"

**STRATEGIC INSIGHTS:**
• "Give me leadership insights on team performance"
• "What are the key risks I should be aware of?"
• "Recommend resource allocation strategies"

I focus on actionable insights that drive business decisions. What would you like to explore?"""
    
    elif any(word in message for word in ['thank', 'thanks']):
        return "You're welcome! I'm here to help you make data-driven leadership decisions. What else would you like to analyze?"
    
    else:
        return f"""I'm your Leadership Quality Assistant, focused on providing strategic insights from your project data. I can help analyze team performance, project health, quality metrics, and resource allocation. 

For '{message}', I'd need more context about what specific leadership insights you're looking for. Are you interested in team performance, project status, quality analysis, or strategic recommendations?"""

# AI Insight Generation Functions
async def generate_velocity_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate velocity and sprint insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    current_sprint = analytics.get("current_sprint")
    
    # Calculate velocity metrics
    total_stories = summary["total_stories"]
    total_defects = summary["total_defects"]
    defect_ratio = (total_defects / total_stories * 100) if total_stories > 0 else 0
    
    insights = []
    recommendations = []
    
    # Velocity analysis
    if current_sprint:
        insights.append(f"**Current Sprint:** {current_sprint.get('name', 'Unknown')}")
        insights.append(f"**Sprint State:** {current_sprint.get('state', 'Unknown')}")
    
    insights.append(f"**Total Stories:** {total_stories}")
    insights.append(f"**Total Defects:** {total_defects}")
    insights.append(f"**Defect Ratio:** {defect_ratio:.1f}%")
    
    # Recommendations based on data
    if defect_ratio > 20:
        recommendations.append("🔴 **High Defect Ratio:** Consider improving testing processes and code review practices")
    elif defect_ratio < 5:
        recommendations.append("🟢 **Low Defect Ratio:** Excellent quality! Maintain current practices")
    else:
        recommendations.append("🟡 **Moderate Defect Ratio:** Room for improvement in quality processes")
    
    if total_stories > 100:
        recommendations.append("📈 **High Story Count:** Consider breaking down large stories for better estimation")
    
    return {
        "success": True,
        "type": "velocity",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_stories": total_stories,
            "total_defects": total_defects,
            "defect_ratio": defect_ratio,
            "total_projects": summary["total_projects"]
        }
    }

async def generate_team_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate team performance insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # Team distribution analysis
    total_assignees = summary["total_assignees"]
    total_issues = summary["total_issues"]
    avg_issues_per_person = total_issues / total_assignees if total_assignees > 0 else 0
    
    insights.append(f"**Total Team Members:** {total_assignees}")
    insights.append(f"**Total Issues:** {total_issues}")
    insights.append(f"**Average Issues per Person:** {avg_issues_per_person:.1f}")
    
    # Project distribution
    project_count = len(projects)
    insights.append(f"**Active Projects:** {project_count}")
    
    # Find most active project
    if projects:
        most_active = max(projects.items(), key=lambda x: x[1]['total_issues'])
        insights.append(f"**Most Active Project:** {most_active[0]} ({most_active[1]['total_issues']} issues)")
    
    # Recommendations
    if avg_issues_per_person > 20:
        recommendations.append("⚠️ **High Workload:** Consider redistributing tasks or adding team members")
    elif avg_issues_per_person < 5:
        recommendations.append("📊 **Low Workload:** Team has capacity for additional work")
    
    if project_count > 5:
        recommendations.append("🎯 **Multiple Projects:** Consider consolidating or prioritizing projects")
    
    return {
        "success": True,
        "type": "team_performance",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_assignees": total_assignees,
            "total_issues": total_issues,
            "avg_issues_per_person": avg_issues_per_person,
            "project_count": project_count
        }
    }

async def generate_project_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate project health insights"""
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # Project health analysis
    healthy_projects = 0
    for project_key, project_data in projects.items():
        stories = project_data['stories']
        defects = project_data['defects']
        defect_ratio = (defects / stories * 100) if stories > 0 else 0
        
        if defect_ratio < 15 and stories > 0:
            healthy_projects += 1
        
        insights.append(f"**{project_key}:** {stories} stories, {defects} defects ({defect_ratio:.1f}% defect ratio)")
    
    health_percentage = (healthy_projects / len(projects) * 100) if projects else 0
    insights.append(f"**Project Health:** {health_percentage:.1f}% of projects are healthy")
    
    # Recommendations
    if health_percentage < 50:
        recommendations.append("🔴 **Low Project Health:** Focus on improving quality processes across projects")
    elif health_percentage > 80:
        recommendations.append("🟢 **Excellent Project Health:** Maintain current practices")
    else:
        recommendations.append("🟡 **Good Project Health:** Continue current practices with minor improvements")
    
    return {
        "success": True,
        "type": "project_health",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_projects": len(projects),
            "healthy_projects": healthy_projects,
            "health_percentage": health_percentage
        }
    }

async def generate_general_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate general insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # General overview
    insights.append(f"**Total Projects:** {summary['total_projects']}")
    insights.append(f"**Total Stories:** {summary['total_stories']}")
    insights.append(f"**Total Defects:** {summary['total_defects']}")
    insights.append(f"**Total Tasks:** {summary['total_tasks']}")
    insights.append(f"**Team Size:** {summary['total_assignees']} members")
    
    # Key recommendations
    recommendations.append("📊 **Data-Driven Decisions:** Use these metrics to guide sprint planning")
    recommendations.append("🎯 **Focus Areas:** Prioritize projects with high defect ratios")
    recommendations.append("👥 **Team Balance:** Ensure even distribution of work across team members")
    recommendations.append("📈 **Continuous Improvement:** Track these metrics over time for trend analysis")
    
    return {
        "success": True,
        "type": "general",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": summary
    }

# API Routes
@app.get("/")
async def root():
    return {"message": "Integration Hub API", "status": "running"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

@app.post("/api/confluence/configure")
async def configure_confluence(config: ConfluenceConfigRequest):
    """Configure Confluence connection"""
    try:
        conf_cfg = ConfluenceConfig(base_url=config.url, email=config.email, api_token=config.api_token)
        client = ConfluenceClient(conf_cfg)
        await client.initialize()
        # quick ping by doing a trivial search
        try:
            _ = await client.search("home", limit=1)
        except Exception:
            pass

        app_state.confluence_configured = True
        app_state.confluence_client = client
        app_state.confluence_config = conf_cfg
        return {"success": True, "message": "Confluence configured"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/confluence/disconnect")
async def disconnect_confluence():
    """Disconnect Confluence connection"""
    try:
        if app_state.confluence_client:
            try:
                await app_state.confluence_client.close()
            except Exception:
                pass
        app_state.confluence_configured = False
        app_state.confluence_client = None
        app_state.confluence_config = None
        return {"success": True, "message": "Confluence disconnected successfully"}
    except Exception as e:
        logger.error(f"Confluence disconnection failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/configure")
async def configure_jira(config: JiraConfigRequest):
    """Configure Jira connection"""
    try:
        # Create JiraConfig object
        jira_config = JiraConfig(
            base_url=config.url,
            email=config.email,
            api_token=config.api_token,
            board_id=config.board_id or ""
        )
        
        # Create Jira client
        jira_client = JiraClient(jira_config)
        
        # Initialize the async client
        await jira_client.initialize()
        
        # Test the connection by trying to get current sprint
        try:
            current_sprint = await jira_client.get_current_sprint()
            logger.info(f"Successfully connected to Jira. Current sprint: {current_sprint}")
        except Exception as e:
            logger.warning(f"Could not get current sprint, but connection may still work: {e}")
        
        # Fetch and cache projects for this session
        try:
            projects = await jira_client.get_projects()
            key_to_name = {}
            keys = []
            for p in projects:
                key = (p.get('key') or '').upper()
                name = p.get('name') or ''
                if key:
                    key_to_name[key] = name
                    keys.append(key)
            app_state.jira_projects = key_to_name
            app_state.jira_project_keys = keys
            # Provide known project keys to NLU for better detection
            try:
                jql_loader.slot_nlu.set_known_projects(keys)
            except Exception as e:
                logger.warning(f"Could not set known project keys in NLU: {e}")
        except Exception as e:
            logger.warning(f"Failed to prefetch Jira projects: {e}")
        
        # Store configuration
        app_state.jira_configured = True
        app_state.jira_client = jira_client
        app_state.jira_config = jira_config
        app_state.jira_board_id = config.board_id or None
        
        return {
            "success": True,
            "message": "Jira configured successfully",
            "config": {
                "url": config.url,
                "email": config.email,
                "board_id": app_state.jira_board_id
            },
            "projects": {
                "keys": app_state.jira_project_keys,
                "map": app_state.jira_projects
            }
        }
    except Exception as e:
        logger.error(f"Jira configuration failed: {e}")
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/jira/disconnect")
async def disconnect_jira():
    """Disconnect Jira connection"""
    try:
        # Close the async client if it exists
        if app_state.jira_client:
            await app_state.jira_client.close()
        
        # Clear configuration
        app_state.jira_configured = False
        app_state.jira_client = None
        app_state.jira_config = None
        app_state.jira_board_id = None
        
        return {
            "success": True,
            "message": "Jira disconnected successfully"
        }
    except Exception as e:
        logger.error(f"Jira disconnection failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/status", tags=["JIRA"], summary="Get JIRA Connection Status")
async def get_jira_status():
    """Get Jira connection status"""
    return {
        "configured": app_state.jira_configured,
        "board_id": app_state.jira_board_id,
        "projects": {
            "keys": app_state.jira_project_keys,
            "map": app_state.jira_projects
        } if app_state.jira_configured else None,
        "config": {
            "url": app_state.jira_config.base_url if app_state.jira_config else None,
            "email": mask_email(app_state.jira_config.email) if app_state.jira_config else None
        } if app_state.jira_config else None
    }

@app.get("/api/confluence/status")
async def get_confluence_status():
    return {
        "configured": app_state.confluence_configured,
        "config": {
            "url": app_state.confluence_config.base_url if app_state.confluence_config else None,
            "email": mask_email(app_state.confluence_config.email) if app_state.confluence_config else None
        } if app_state.confluence_config else None
    }

@app.get("/api/jira/projects", tags=["JIRA"], summary="Get cached Jira projects")
async def get_cached_projects():
    """Return cached Jira project keys and names for this session"""
    if not app_state.jira_configured:
        raise HTTPException(status_code=400, detail="Jira not configured")
    # If cache empty but client exists, try to refresh once
    if not app_state.jira_projects and app_state.jira_client:
        try:
            projects = await app_state.jira_client.get_projects()
            key_to_name = {}
            keys = []
            for p in projects:
                key = (p.get('key') or '').upper()
                name = p.get('name') or ''
                if key:
                    key_to_name[key] = name
                    keys.append(key)
            app_state.jira_projects = key_to_name
            app_state.jira_project_keys = keys
            try:
                jql_loader.slot_nlu.set_known_projects(keys)
            except Exception:
                pass
        except Exception as e:
            logger.error(f"Failed to fetch Jira projects: {e}")
            raise HTTPException(status_code=500, detail="Failed to fetch Jira projects")
    return {
        "success": True,
        "projects": {
            "keys": app_state.jira_project_keys,
            "map": app_state.jira_projects
        }
    }

@app.get("/api/jira/sprint/current")
async def get_current_sprint():
    """Get current sprint information"""
    # Require at least one integration (Jira or Confluence) to be configured
    jira_ok = app_state.jira_configured and app_state.jira_client is not None
    conf_ok = app_state.confluence_configured and app_state.confluence_client is not None
    if not (jira_ok or conf_ok):
        raise HTTPException(status_code=400, detail="No integrations configured. Please configure Jira or Confluence.")
    
    try:
        current_sprint = await app_state.jira_client.get_current_sprint()
        return {
            "success": True,
            "sprint": current_sprint
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/search")
async def search_jira_issues(request: Dict[str, Any]):
    """Search Jira issues"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        jql = request.get("jql", "project is not EMPTY")
        max_results = request.get("max_results", 50)
        fields = request.get("fields", None)
        
        result = await app_state.jira_client.search(jql, max_results)
        return {
            "success": True,
            "result": result
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/analytics")
async def get_jira_analytics():
    """Get comprehensive Jira analytics with AI insights"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Initialize analytics engine if not already done
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = AdvancedAIEngine(app_state.jira_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Generate comprehensive analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        return {
            "success": True,
            "analytics": analytics
        }
    except Exception as e:
        logger.error(f"Advanced analytics error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/predictive-analysis")
async def get_predictive_analysis(request: Dict[str, Any]):
    """Get predictive analysis and forecasting"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        query = request.get("query", "What are the trends and predictions for our team performance?")
        
        # Initialize AI components if not already done
        if not app_state.ai_engine:
            app_state.ai_engine = AdvancedAIEngine(app_state.jira_client)
        if not app_state.analytics_engine:
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get historical data for prediction
        historical_jql = "project is not EMPTY AND updated >= -90d ORDER BY updated DESC"
        historical_data = await app_state.jira_client.search(historical_jql, max_results=1000)
        
        # Generate prediction
        prediction = app_state.ai_engine.generate_predictive_analysis(query, historical_data)
        
        return {
            "success": True,
            "prediction": prediction,
            "data_points": len(historical_data.get('issues', [])),
            "timeframe": "next_2_weeks"
        }
    except Exception as e:
        logger.error(f"Predictive analysis error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/anomaly-detection")
async def get_anomaly_detection():
    """Get anomaly detection results"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Initialize analytics engine if not already done
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = AdvancedAIEngine(app_state.jira_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get current analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        # Detect anomalies
        anomalies = app_state.analytics_engine.detect_anomalies(analytics)
        
        return {
            "success": True,
            "anomalies": [anomaly.__dict__ for anomaly in anomalies],
            "total_anomalies": len(anomalies),
            "critical_count": len([a for a in anomalies if a.severity == 'critical']),
            "high_count": len([a for a in anomalies if a.severity == 'high'])
        }
    except Exception as e:
        logger.error(f"Anomaly detection error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/intelligent-recommendations")
async def get_intelligent_recommendations(request: Dict[str, Any]):
    """Get AI-powered intelligent recommendations"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        query = request.get("query", "What recommendations do you have for improving our team performance?")
        
        # Initialize AI components if not already done
        if not app_state.ai_engine:
            app_state.ai_engine = AdvancedAIEngine(app_state.jira_client)
        if not app_state.analytics_engine:
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get comprehensive analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        # Generate intelligent response with recommendations
        response = app_state.ai_engine.generate_intelligent_response(query, analytics)
        
        return {
            "success": True,
            "recommendations": response,
            "analytics_summary": analytics.get('summary', {}),
            "insights_count": len(analytics.get('ai_insights', [])),
            "anomalies_count": len(analytics.get('anomalies', []))
        }
    except Exception as e:
        logger.error(f"Intelligent recommendations error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/ai-insights")
async def get_ai_insights(request: Dict[str, Any]):
    """Get AI-powered insights and recommendations"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        insight_type = request.get("type", "general")
        jira_client = app_state.jira_client
        
        # Get analytics data directly from analytics engine
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = AdvancedAIEngine(app_state.jira_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        if insight_type == "velocity":
            return await generate_velocity_insights(analytics, jira_client)
        elif insight_type == "team_performance":
            return await generate_team_insights(analytics, jira_client)
        elif insight_type == "project_health":
            return await generate_project_insights(analytics, jira_client)
        else:
            return await generate_general_insights(analytics, jira_client)
            
    except Exception as e:
        logger.error(f"AI insights error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/export")
async def export_jira_analytics(request: Dict[str, Any]):
    """Export Jira analytics to various formats"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        export_format = request.get("format", "json")
        
        # Get analytics data
        analytics_response = await get_jira_analytics()
        analytics = analytics_response["analytics"]
        
        if export_format == "json":
            return {
                "success": True,
                "data": analytics,
                "format": "json",
                "filename": f"jira_analytics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
            }
        elif export_format == "csv":
            # Convert to CSV format
            csv_data = convert_analytics_to_csv(analytics)
            return {
                "success": True,
                "data": csv_data,
                "format": "csv",
                "filename": f"jira_analytics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
            }
        else:
            raise HTTPException(status_code=400, detail="Unsupported export format")
            
    except Exception as e:
        logger.error(f"Export error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

def convert_analytics_to_csv(analytics: Dict[str, Any]) -> str:
    """Convert analytics data to CSV format with robust error handling"""
    import csv
    import io
    
    output = io.StringIO()
    writer = csv.writer(output)
    
    # Write summary with safe access
    writer.writerow(["Metric", "Value"])
    summary = analytics.get("summary", {})
    for key, value in summary.items():
        writer.writerow([key.replace("_", " ").title(), value])
    
    writer.writerow([])  # Empty row
    
    # Write project details with safe access
    writer.writerow(["Project", "Stories", "Defects", "Tasks", "Total Issues", "Assignees"])
    projects = analytics.get("projects", {})
    for project_key, project_data in projects.items():
        writer.writerow([
            project_key,
            project_data.get("stories", 0),
            project_data.get("defects", 0),
            project_data.get("tasks", 0),
            project_data.get("total_issues", 0),
            project_data.get("assignee_count", 0)
        ])
    
    return output.getvalue()

@app.post("/api/chat", tags=["Chat"], summary="Chat with AI Assistant")
async def chat_endpoint(request: ChatRequest):
    """Handle chat messages with ticket key detection, training-first approach, then AI fallback"""
    jira_available = bool(app_state.jira_configured and app_state.jira_client)

    user_q = request.message.strip()
    q_lower_for_global = user_q.lower()
    
    # Quick global pagination: handle "show me others" using last stored JQL
    if jira_available:
        page_keywords = ['show me others', 'show others', 'show more', 'next page', 'more items', 'remaining']
        if any(k in q_lower_for_global for k in page_keywords) and getattr(app_state, 'last_list_jql', None):
            try:
                jql = app_state.last_list_jql or ""
                total_count = app_state.last_list_total or await app_state.jira_client.count(jql)
                next_offset = (app_state.last_list_offset or 0) + 10
                if next_offset >= max(total_count, 0):
                    return {"success": True, "response": f"That's all for now. Already showed {total_count} items."}
                result = await app_state.jira_client.search(jql, max_results=10, start_at=next_offset, fields=["key", "summary", "status", "priority", "assignee"])
                issues = result.get("issues", [])
                if not issues:
                    return {"success": True, "response": "No more items found."}
                start_num = next_offset + 1
                lines = [f"Found {total_count} total items. Showing items {start_num}-{start_num + len(issues) - 1}:"]
                for i, issue in enumerate(issues, start_num):
                    key = issue.get("key", "Unknown")
                    fields = issue.get("fields", {})
                    summary = fields.get("summary", "No summary")
                    status = fields.get("status", {}).get("name", "Unknown")
                    priority = fields.get("priority", {}).get("name", "Medium")
                    assignee = fields.get("assignee", {})
                    assignee_name = assignee.get("displayName", "Unassigned") if assignee else "Unassigned"
                    lines.append(f"{i}. {key}: {summary}")
                    lines.append(f"   Status: {status} | Priority: {priority} | Assignee: {assignee_name}")
                    lines.append("")
                # Inline tips removed; export handled by UI buttons
                # Leadership summary overlay
                try:
                    summary = ai_summarize(user_q, issues, 'list')
                except Exception:
                    summary = None
                # Update listing context
                app_state.last_list_offset = next_offset
                app_state.last_list_total = total_count
                resp_text = (summary + "\n\n" if summary else "") + "\n".join(lines) + f"\n\n[debug] JQL: {jql}"
                return {"success": True, "mode": "training", "jql": jql, "raw": {"count": len(issues), "total_count": total_count, "issues": issues, "start_at": next_offset}, "response": resp_text}
            except Exception as e:
                logger.error(f"Pagination using last JQL failed: {e}")
                # fall through to normal handling

        # Global export: allow export right after any list without requiring re-match
        export_keywords = ['export to excel', 'download excel', 'export excel', 'excel file', 'download list']
        if any(k in q_lower_for_global for k in export_keywords):
            try:
                export_jql = getattr(app_state, 'last_list_jql', '')
                if not export_jql:
                    return {"success": True, "mode": "export", "response": "Tell me a project (e.g., CCM) or a person (e.g., Ashwin)."}
                total_count = await app_state.jira_client.count(export_jql)
                if total_count <= 0:
                    return {"success": True, "mode": "export", "response": "No items found to export."}
                result = await app_state.jira_client.search(export_jql, max_results=min(total_count, 1000), fields=["key", "summary", "status", "priority", "assignee", "created", "updated"])
                issues = result.get("issues", [])
                import pandas as pd
                from io import BytesIO
                import base64
                data = []
                for issue in issues:
                    assignee = issue.get("fields", {}).get("assignee", {})
                    assignee_name = assignee.get("displayName", "Unassigned") if assignee else "Unassigned"
                    data.append({
                        "Ticket": issue.get("key", ""),
                        "Summary": issue.get("fields", {}).get("summary", ""),
                        "Status": issue.get("fields", {}).get("status", {}).get("name", ""),
                        "Priority": issue.get("fields", {}).get("priority", {}).get("name", ""),
                        "Assignee": assignee_name,
                        "Created": issue.get("fields", {}).get("created", ""),
                        "Updated": issue.get("fields", {}).get("updated", "")
                    })
                df = pd.DataFrame(data)
                excel_buffer = BytesIO()
                df.to_excel(excel_buffer, index=False, sheet_name="Jira Items")
                excel_buffer.seek(0)
                excel_data = base64.b64encode(excel_buffer.getvalue()).decode()
                filename = f"jira_items_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                return {
                    "success": True,
                    "mode": "export",
                    "response": f"📊 Excel file ready! Downloaded {len(issues)} items as '{filename}'",
                    "export_data": {
                        "filename": filename,
                        "data": excel_data,
                        "count": len(issues),
                        "total_count": total_count
                    }
                }
            except Exception as e:
                logger.error(f"Global export failed: {e}")
                return {"success": True, "mode": "export", "response": f"Export failed: {str(e)}"}

    # Check for ticket key in the query (case-insensitive)
    ticket_match = re.search(ISSUE_KEY_RE, user_q, re.IGNORECASE)
    if ticket_match and jira_available:
        ticket_key = ticket_match.group(1)
        try:
            # Get direct issue details
            issue = await app_state.jira_client.get_issue(ticket_key, fields=[
                "summary", "status", "issuetype", "priority", "assignee", "reporter", "project"
            ])
            
            # Extract key information
            fields = issue.get("fields", {})
            summary = fields.get("summary", "No summary")
            status = fields.get("status", {}).get("name", "Unknown")
            issue_type = fields.get("issuetype", {}).get("name", "Unknown")
            priority = fields.get("priority", {}).get("name", "Unknown")
            assignee = fields.get("assignee", {})
            assignee_name = assignee.get("displayName", "Unassigned") if assignee else "Unassigned"
            reporter = fields.get("reporter", {}) or {}
            reporter_name = reporter.get("displayName", "Unknown")
            reporter_email = reporter.get("emailAddress") or ""
            project = fields.get("project", {}).get("key", "Unknown")
            
            # Create grounded summary with structured formatting and a brief leadership note
            canonical_key = issue.get("key", ticket_key).upper()
            response_text = (
                f"{canonical_key}: {summary}\n\n"
                f"- Status: {status}\n"
                f"- Type: {issue_type}\n"
                f"- Priority: {priority}\n"
                f"- Assignee: {assignee_name}\n"
                f"- Reporter: {reporter_name}{f' ({reporter_email})' if reporter_email else ''}\n"
                f"- Project: {project}\n\n"
                f"Leadership note: {assignee_name} owns this {issue_type.lower()} currently {status.lower()}. Priority is {priority.lower()}."
            )
            
            return {
                "success": True,
                "response": response_text
            }
            
        except Exception as e:
            logger.error(f"Failed to get issue {ticket_key}: {e}")
            return {
                "success": True,
                "response": f"Sorry, I couldn't find details for {ticket_key}. The ticket might not exist or there might be a connection issue."
            }
    
    # If Confluence is configured and the user asks about wiki/Confluence/docs, try Confluence first
    if app_state.confluence_configured and app_state.confluence_client:
        try:
            # Heuristic title extraction for phrases like "give me X from confluence"
            q_lower = user_q.lower()
            cleaned = q_lower
            for token in [
                'give me', 'show me', 'get me', 'please', 'from confluence', 'in confluence',
                'from wiki', 'in wiki', 'confluence', 'wiki', 'article', 'page'
            ]:
                cleaned = cleaned.replace(token, ' ')
            cleaned = ' '.join(cleaned.split()).strip()

            # Try direct title search first
            title_query = cleaned or user_q
            pages = await app_state.confluence_client.search_pages(title_query, limit=1)
            if pages:
                p = pages[0]
                content = p.get("content") or {}
                content_id = content.get("id") or p.get("id")
                full = await app_state.confluence_client.get_page(content_id)
                if full:
                    title = full.get("title", "Untitled")
                    storage = full.get("body", {}).get("storage", {}).get("value", "")
                    text = ConfluenceClient.storage_html_to_text(storage)
                    return {"success": True, "response": f"{title}\n\n{text}"}

            # Fallback: general text search with excerpts
            results = await app_state.confluence_client.search(title_query, limit=5)
            if results:
                lines = ["Here are relevant Confluence pages:"]
                for r in results:
                    title = r.get("title") or r.get("content", {}).get("title", "Untitled")
                    content = r.get("content") or {}
                    content_id = content.get("id") or r.get("id")
                    excerpt = (r.get("excerpt") or "").replace("<em>", "").replace("</em>", "")
                    lines.append(f"- {title} (ID: {content_id})\n  {excerpt[:180]}...")
                return {"success": True, "response": "\n".join(lines)}
        except Exception as e:
            logger.warning(f"Confluence search failed: {e}")

    # No ticket key detected, proceed with enhanced slot-based flow
    # New AI intent router first
    try:
        ai_answer = await ai_handle(user_q, app_state.jira_client)
        if ai_answer and isinstance(ai_answer, dict) and ai_answer.get("text") and not ai_answer.get("text").startswith("I couldn’t map"):
            # Persist list context if it's a listing intent
            if 'list' in (jql_loader.intent_of(user_q) or '').lower() or 'stories' in user_q.lower() or 'list' in user_q.lower():
                app_state.last_list_jql = ai_answer.get("jql") or app_state.last_list_jql
                app_state.last_list_offset = 0
            return {"success": True, "mode": "ai_router", "response": ai_answer.get("text"), "jql": ai_answer.get("jql")}
    except Exception as _e:
        logger.info(f"AI router fallback: {_e}")

    match = jql_loader.find(user_q)
    if match and jira_available:
        jql = match["jql"]

        # Extract slots from the user query for friendlier responses and JQL fixes
        try:
            slots = jql_loader.slot_nlu.extract_slots(user_q)
        except Exception:
            slots = None

        # Resolve assignee to accountId on Jira Cloud to avoid name-based failures
        try:
            if slots and getattr(slots, "assignee", None) and slots.assignee and slots.assignee.value:
                assignee_clause = await app_state.jira_client.build_assignee_jql(slots.assignee.value)
                # If JQL already has an assignee clause, replace it; else append
                import re as _re
                if _re.search(r"assignee\s*[=~]", jql, _re.IGNORECASE):
                    jql = _re.sub(r"assignee\s*[=~]\s*\S+", assignee_clause, jql, flags=_re.IGNORECASE)
                else:
                    connector = " AND " if jql.strip() else ""
                    jql = f"{jql}{connector}{assignee_clause}"
        except Exception:
            pass

        # Update context if project was detected
        if "slots_used" in match and "project" in match["slots_used"]:
            jql_loader.set_context(project=match["slots_used"]["project"])
        
        # Special handling for project details queries - provide segmented breakdown
        if match.get("intent") == "project_general" and slots and getattr(slots, "project", None) and slots.project and slots.project.value:
            project_key = slots.project.value
            try:
                # Get segmented breakdown for the project
                stories_jql = f"project = \"{project_key}\" AND issuetype = Story"
                bugs_jql = f"project = \"{project_key}\" AND issuetype in (Bug, Defect)"
                tasks_jql = f"project = \"{project_key}\" AND issuetype = Task"
                total_jql = f"project = \"{project_key}\""
                
                stories_count = await app_state.jira_client.count(stories_jql)
                bugs_count = await app_state.jira_client.count(bugs_jql)
                tasks_count = await app_state.jira_client.count(tasks_jql)
                total_count = await app_state.jira_client.count(total_jql)
                
                response_text = (
                    f"**{project_key} Project Breakdown:**\n\n"
                    f"- Stories: {stories_count}\n"
                    f"- Defects: {bugs_count}\n"
                    f"- Tasks: {tasks_count}\n"
                    f"- Total Issues: {total_count}\n\n"
                    f"Leadership note: {project_key} has {total_count} total items with {stories_count} stories, {bugs_count} defects, and {tasks_count} tasks."
                )
                
                return {
                    "success": True,
                    "mode": "training",
                    "jql": jql,
                    "raw": {"count": total_count, "stories": stories_count, "bugs": bugs_count, "tasks": tasks_count},
                    "response": response_text + f"\n\n[debug] JQL: {jql}",
                    "slots_used": match.get("slots_used", {})
                }
            except Exception as e:
                logger.error(f"Error getting project breakdown for {project_key}: {e}")
                # Fall back to regular count
                pass
        
        # Check if user wants to export to Excel
        export_keywords = ['export to excel', 'download excel', 'export excel', 'excel file', 'download list']
        wants_export = any(keyword in user_q.lower() for keyword in export_keywords)
        
        if wants_export:
            try:
                # Prefer current JQL; if not present, reuse the last listing JQL
                export_jql = jql or getattr(app_state, 'last_list_jql', '')
                if not export_jql:
                    return {"success": True, "mode": "export", "response": "Tell me a project (e.g., CCM) or a person (e.g., Ashwin)."}
                # Get all items for export
                total_count = await app_state.jira_client.count(export_jql)
                if total_count > 0:
                    # Get all items (up to 1000 for performance)
                    result = await app_state.jira_client.search(export_jql, max_results=min(total_count, 1000), fields=["key", "summary", "status", "priority", "assignee", "created", "updated"])
                    issues = result.get("issues", [])
                    
                    # Create Excel file
                    import pandas as pd
                    from io import BytesIO
                    import base64
                    
                    # Prepare data for Excel
                    data = []
                    for issue in issues:
                        assignee = issue.get("fields", {}).get("assignee", {})
                        assignee_name = assignee.get("displayName", "Unassigned") if assignee else "Unassigned"
                        
                        data.append({
                            "Ticket": issue.get("key", ""),
                            "Summary": issue.get("fields", {}).get("summary", ""),
                            "Status": issue.get("fields", {}).get("status", {}).get("name", ""),
                            "Priority": issue.get("fields", {}).get("priority", {}).get("name", ""),
                            "Assignee": assignee_name,
                            "Created": issue.get("fields", {}).get("created", ""),
                            "Updated": issue.get("fields", {}).get("updated", "")
                        })
                    
                    # Create DataFrame and export to Excel
                    df = pd.DataFrame(data)
                    excel_buffer = BytesIO()
                    df.to_excel(excel_buffer, index=False, sheet_name="Jira Items")
                    excel_buffer.seek(0)
                    
                    # Encode to base64 for download
                    excel_data = base64.b64encode(excel_buffer.getvalue()).decode()
                    filename = f"jira_items_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                    
                    return {
                        "success": True,
                        "mode": "export",
                        "response": f"📊 Excel file ready! Downloaded {len(issues)} items as '{filename}'",
                        "export_data": {
                            "filename": filename,
                            "data": excel_data,
                            "count": len(issues),
                            "total_count": total_count
                        },
                        "slots_used": match.get("slots_used", {}) if match else {}
                    }
                else:
                    return {
                        "success": True,
                        "mode": "export",
                        "response": "No items found to export.",
                        "slots_used": match.get("slots_used", {}) if match else {}
                    }
            except Exception as e:
                logger.error(f"Error exporting to Excel: {e}")
                return {
                    "success": True,
                    "mode": "export",
                    "response": f"Export failed: {str(e)}",
                    "slots_used": match.get("slots_used", {}) if match else {}
                }
        
        # Check if user wants to list actual items instead of just count
        list_keywords = ['list', 'show', 'give me', 'display', 'what are', 'which are', 'what is', 'which is', 'what does', 'which does', 'what work', 'which work', 'what tasks', 'which tasks', 'what stories', 'which stories', 'what bugs', 'which bugs', 'what issues', 'which issues']
        wants_list = any(keyword in user_q.lower() for keyword in list_keywords)
        
        if wants_list:
            try:
                # Check if user wants to see more items (pagination)
                page_keywords = ['show me others', 'show others', 'show more', 'next page', 'more items', 'remaining']
                wants_more = any(keyword in user_q.lower() for keyword in page_keywords)
                
                # Get total count first
                total_count = await app_state.jira_client.count(jql)
                
                if wants_more:
                    # For pagination, use stored offset to avoid repeating first page
                    next_offset = (app_state.last_list_offset or 0) + 10
                    result = await app_state.jira_client.search(jql, max_results=10, start_at=next_offset, fields=["key", "summary", "status", "priority", "assignee"])
                    start_num = next_offset + 1
                else:
                    # First page - get first 10 items
                    result = await app_state.jira_client.search(jql, max_results=10, fields=["key", "summary", "status", "priority", "assignee"])
                    start_num = 1
                
                issues = result.get("issues", [])
                
                if issues:
                    # Build numbered list with key details
                    lines = [f"Found {total_count} total items. Showing items {start_num}-{start_num + len(issues) - 1}:"]
                    for i, issue in enumerate(issues, start_num):
                        key = issue.get("key", "Unknown")
                        summary = issue.get("fields", {}).get("summary", "No summary")
                        status = issue.get("fields", {}).get("status", {}).get("name", "Unknown")
                        priority = issue.get("fields", {}).get("priority", {}).get("name", "Medium")
                        assignee = issue.get("fields", {}).get("assignee", {})
                        assignee_name = assignee.get("displayName", "Unassigned") if assignee else "Unassigned"
                        
                        lines.append(f"{i}. {key}: {summary}")
                        lines.append(f"   Status: {status} | Priority: {priority} | Assignee: {assignee_name}")
                        lines.append("")  # Empty line for readability
                    
                    # Tips removed; UI has export button
                    
                    # Leadership summary overlay
                    try:
                        summary = ai_summarize(user_q, issues, 'list')
                    except Exception:
                        summary = None
                    response_text = ((summary + "\n\n") if summary else "") + "\n".join(lines) + f"\n\n[debug] JQL: {jql}"
                    # Save context for pagination/export continuity
                    app_state.last_list_jql = jql
                    app_state.last_list_offset = 0 if not wants_more else (start_num - 1)
                    app_state.last_list_total = total_count
                    
                    return {
                        "success": True,
                        "mode": "training",
                        "jql": jql,
                        "raw": {"count": len(issues), "total_count": total_count, "issues": issues, "start_at": start_num - 1},
                        "response": response_text,
                        "slots_used": match.get("slots_used", {})
                    }
                else:
                    return {
                        "success": True,
                        "mode": "training",
                        "jql": jql,
                        "raw": {"count": 0, "total_count": total_count},
                        "response": "No items found matching your criteria.\n\n[debug] JQL: {jql}",
                        "slots_used": match.get("slots_used", {})
                    }
            except Exception as e:
                logger.error(f"Error listing items: {e}")
                # Fall back to count
                pass
        
        # Use count() for fast counting, fallback to search().total
        try:
            count = await app_state.jira_client.count(jql)
        except Exception:
            # Fallback if approximate-count not available
            result = await app_state.jira_client.search(jql, max_results=0)
            count = result.get("total", 0)

        # Build a friendlier natural response summarizing applied filters
        fragments = []
        if slots and getattr(slots, "issuetype", None) and slots.issuetype and slots.issuetype.value:
            fragments.append(f"{slots.issuetype.value}s")
        else:
            fragments.append("issues")
        if slots and getattr(slots, "status", None) and slots.status and slots.status.value:
            fragments.insert(0, slots.status.value)
        if slots and getattr(slots, "assignee", None) and slots.assignee and slots.assignee.value:
            fragments.append(f"for {slots.assignee.value}")
        if slots and getattr(slots, "project", None) and slots.project and slots.project.value:
            fragments.append(f"in {slots.project.value}")
        natural = f"Found {count} {' '.join(fragments)}.\n\n[debug] JQL: {jql}"
        
        return {
            "success": True,
            "mode": "training",
            "jql": jql,
            "raw": {"count": count},
            "response": natural,
            "slots_used": match.get("slots_used", {})
        }

    # Fallback: Use advanced JQL generator with fuzzy matching and project detection
    if not jira_available:
        if app_state.confluence_configured and app_state.confluence_client:
            try:
                q_lower = user_q.lower()
                cleaned = q_lower
                for token in [
                    'give me', 'show me', 'get me', 'please', 'from confluence', 'in confluence',
                    'from wiki', 'in wiki', 'confluence', 'wiki', 'article', 'page'
                ]:
                    cleaned = cleaned.replace(token, ' ')
                cleaned = ' '.join(cleaned.split()).strip() or user_q

                pages = await app_state.confluence_client.search_pages(cleaned, limit=1)
                if pages:
                    p = pages[0]
                    content = p.get("content") or {}
                    content_id = content.get("id") or p.get("id")
                    full = await app_state.confluence_client.get_page(content_id)
                    if full:
                        title = full.get("title", "Untitled")
                        storage = full.get("body", {}).get("storage", {}).get("value", "")
                        text = ConfluenceClient.storage_html_to_text(storage)
                        return {"success": True, "response": f"{title}\n\n{text}"}

                results = await app_state.confluence_client.search(user_q, limit=5)
                if results:
                    lines = ["Here are relevant Confluence pages:"]
                    for r in results:
                        title = r.get("title") or r.get("content", {}).get("title", "Untitled")
                        content = r.get("content") or {}
                        content_id = content.get("id") or r.get("id")
                        excerpt = (r.get("excerpt") or "").replace("<em>", "").replace("</em>", "")
                        lines.append(f"- {title} (ID: {content_id})\n  {excerpt[:180]}...")
                    return {"success": True, "response": "\n".join(lines)}

                return {"success": True, "response": f"I couldn't find Confluence results for: '{user_q}'. Try a different title or keyword."}
            except Exception as e:
                logger.warning(f"Confluence fallback search failed: {e}")
                return {"success": False, "response": f"Confluence search failed: {e}"}

        return {"success": True, "response": "Jira is not configured, and Confluence is unavailable. Please connect one integration to proceed."}

    try:
        from utils.advanced_jql_generator import AdvancedJQLGenerator
        
        # Extract project context from enhanced loader or previous messages
        project_ctx = jql_loader.get_context_project()
        if not project_ctx and app_state.messages:
            # Look for project mentions in recent messages
            recent_messages = app_state.messages[-5:]  # Last 5 messages
            for msg in recent_messages:
                if isinstance(msg, dict) and 'message' in msg:
                    # Simple project detection in recent context
                    project_match = re.search(r'\b([A-Z]{2,5})\b', msg['message'])
                    if project_match:
                        project_ctx = project_match.group(1)
                        jql_loader.set_context(project=project_ctx)
                        break
        
        generator = AdvancedJQLGenerator(app_state.jira_client, project_ctx)
        result = await generator.generate_jql(user_q)
        
        if result['success']:
            return {
                "success": True,
                "mode": "advanced_jql",
                "jql": result.get('jql', ''),
                "raw": {"count": result.get('count', 0)},
                "response": result['response']
            }
        else:
            return {
                "success": True,
                "mode": "advanced_jql_error",
                "response": result['response']
            }
            
    except Exception as e:
        logger.error(f"Advanced JQL generator error: {e}")
        return {
            "success": True,
            "mode": "fallback_error",
            "response": f"Sorry, I encountered an error while processing your query: {str(e)}"
        }

@app.get("/api/context", tags=["Context"], summary="Get Current Context")
async def get_context():
    """Get current project/board/sprint context"""
    return {
        "success": True,
        "context": {
            "project": jql_loader.get_context_project(),
            "board": None,  # TODO: implement board context
            "sprint": None  # TODO: implement sprint context
        }
    }

@app.post("/api/context", tags=["Context"], summary="Set Context")
async def set_context(request: Dict[str, Any]):
    """Set project/board/sprint context"""
    project = request.get("project")
    board = request.get("board")
    sprint = request.get("sprint")
    
    jql_loader.set_context(project=project, board=board, sprint=sprint)
    
    return {
        "success": True,
        "message": "Context updated successfully",
        "context": {
            "project": project,
            "board": board,
            "sprint": sprint
        }
    }

@app.get("/api/jira/diagnostics/search")
async def jira_search_diagnostics():
    """Health check for Enhanced JQL endpoints"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    diagnostics = {
        "enhanced_search_ok": False,
        "approximate_count_ok": False,
        "errors": []
    }
    
    try:
        # Test Enhanced JQL search
        test_jql = "project is not EMPTY"
        search_result = await app_state.jira_client.search(test_jql, max_results=1)
        diagnostics["enhanced_search_ok"] = True
    except Exception as e:
        diagnostics["errors"].append(f"Enhanced search failed: {str(e)}")
    
    try:
        # Test approximate count
        test_jql = "project is not EMPTY"
        count_result = await app_state.jira_client.count(test_jql)
        diagnostics["approximate_count_ok"] = True
    except Exception as e:
        diagnostics["errors"].append(f"Approximate count failed: {str(e)}")
    
    return {"success": True, "data": diagnostics}


@app.get("/api/chat/history")
async def get_chat_history():
    """Get chat history"""
    return {"messages": app_state.messages}

@app.post("/api/chat/clear")
async def clear_chat():
    """Clear chat history"""
    app_state.messages = []
    return {"success": True, "message": "Chat cleared"}

@app.post("/api/chat/enhanced", tags=["Chat"], summary="Enhanced Chat with JQL Processing")
async def enhanced_chat_endpoint(request: ChatRequest):
    """Handle chat messages with enhanced JQL processing"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing enhanced message: '{message}'")
        logger.info(f"Jira configured: {app_state.jira_configured}")
        logger.info(f"Jira client exists: {app_state.jira_client is not None}")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize enhanced JQL processor if needed
        if not app_state.enhanced_jql_processor:
            logger.info("🔍 Creating Enhanced JQL Processor...")
            app_state.enhanced_jql_processor = EnhancedJQLProcessor(app_state.jira_client)
        
        # Process query with enhanced JQL processor
        logger.info("🔍 About to call enhanced_jql_processor.process_query")
        result = await app_state.enhanced_jql_processor.process_query(message, ResponseFormat.TEXT)
        logger.info(f"🔍 Enhanced JQL processor returned: {result}")
        
        # Generate response
        response = result.get('response', 'I apologize, but I encountered an issue processing your request.')
        logger.info(f"🔍 Final enhanced response: {response}")
        
        # Store message in history
        app_state.messages.append({
            "id": len(app_state.messages) + 1,
            "message": message,
            "response": response,
            "timestamp": datetime.now().isoformat(),
            "metadata": {
                "ai_enhanced": True,
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        })
        
        return {
            "response": response,
            "metadata": {
                "ai_enhanced": True,
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        }
        
    except Exception as e:
        logger.error(f"Enhanced chat error: {e}")
        return {
            "response": f"I encountered an error while processing your request: {str(e)}",
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.post("/api/chat/json", tags=["Chat"], summary="Chat with JSON Response Format")
async def json_chat_endpoint(request: ChatRequest):
    """Handle chat messages with JSON response format"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing JSON message: '{message}'")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize enhanced JQL processor if needed
        if not app_state.enhanced_jql_processor:
            logger.info("🔍 Creating Enhanced JQL Processor...")
            app_state.enhanced_jql_processor = EnhancedJQLProcessor(app_state.jira_client)
        
        # Process query with enhanced JQL processor in JSON mode
        logger.info("🔍 About to call enhanced_jql_processor.process_query (JSON mode)")
        result = await app_state.enhanced_jql_processor.process_query(message, ResponseFormat.JSON)
        logger.info(f"🔍 Enhanced JQL processor returned JSON: {result}")
        
        # Parse JSON response
        try:
            response_data = json.loads(result.get('response', '{}'))
        except json.JSONDecodeError:
            response_data = {"error": "Invalid JSON response", "raw_response": result.get('response', '')}
        
        return {
            "response": response_data,
            "metadata": {
                "ai_enhanced": True,
                "format": "json",
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        }
        
    except Exception as e:
        logger.error(f"JSON chat error: {e}")
        return {
            "response": {"error": f"I encountered an error while processing your request: {str(e)}"},
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.get("/api/jira/board/{board_id}/sprint/history", tags=["JIRA"], summary="Get Sprint Velocity History")
async def get_sprint_history(board_id: int):
    """Get last 3 sprints velocity for leadership dashboards"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Get last 3 sprints
        sprints = await app_state.jira_client.get_sprint_history(board_id, limit=3)
        
        sprint_data = []
        for sprint in sprints:
            # Get sprint metrics
            sprint_issues = await app_state.jira_client.search(
                f"sprint = {sprint['id']}", 
                max_results=1000
            )
            
            # Calculate velocity
            completed_issues = [
                issue for issue in sprint_issues.get('issues', [])
                if issue.get('fields', {}).get('status', {}).get('name') in ['Done', 'Closed']
            ]
            
            velocity = len(completed_issues)
            
            sprint_data.append({
                "sprint_id": sprint['id'],
                "sprint_name": sprint['name'],
                "start_date": sprint.get('startDate'),
                "end_date": sprint.get('endDate'),
                "velocity": velocity,
                "total_issues": len(sprint_issues.get('issues', [])),
                "completion_rate": velocity / len(sprint_issues.get('issues', [])) * 100 if sprint_issues.get('issues') else 0
            })
        
        return create_success_response(sprint_data, "Sprint history retrieved successfully")
        
    except Exception as e:
        logger.error(f"Sprint history error: {e}")
        return create_error_response("Sprint history failed", str(e))

@app.get("/api/jira/blockers", tags=["JIRA"], summary="Get Blocked Issues")
async def get_blocked_issues():
    """Show flagged issues or status=Blocked"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Get blocked issues
        blocked_jql = "status = Blocked OR status = Waiting"
        blocked_issues = await app_state.jira_client.search(blocked_jql, max_results=100)
        
        # Get flagged issues
        flagged_jql = "labels = flagged OR priority = Highest"
        flagged_issues = await app_state.jira_client.search(flagged_jql, max_results=100)
        
        # Process blocked issues
        blocked_data = []
        for issue in blocked_issues.get('issues', []):
            blocked_data.append({
                "key": issue['key'],
                "summary": issue['fields']['summary'],
                "status": issue['fields']['status']['name'],
                "assignee": issue['fields'].get('assignee', {}).get('displayName', 'Unassigned'),
                "project": issue['fields']['project']['key'],
                "created": issue['fields']['created'],
                "updated": issue['fields']['updated'],
                "url": f"{app_state.jira_config.base_url}/browse/{issue['key']}"
            })
        
        # Process flagged issues
        flagged_data = []
        for issue in flagged_issues.get('issues', []):
            flagged_data.append({
                "key": issue['key'],
                "summary": issue['fields']['summary'],
                "priority": issue['fields'].get('priority', {}).get('name', 'Medium'),
                "labels": [label for label in issue['fields'].get('labels', [])],
                "assignee": issue['fields'].get('assignee', {}).get('displayName', 'Unassigned'),
                "project": issue['fields']['project']['key'],
                "url": f"{app_state.jira_config.base_url}/browse/{issue['key']}"
            })
        
        return create_success_response({
            "blocked_issues": blocked_data,
            "flagged_issues": flagged_data,
            "blocked_count": len(blocked_data),
            "flagged_count": len(flagged_data)
        }, "Blocked and flagged issues retrieved successfully")
        
    except Exception as e:
        logger.error(f"Blocked issues error: {e}")
        return create_error_response("Blocked issues retrieval failed", str(e))

@app.post("/api/chat/advanced", tags=["Chat"], summary="Advanced Chat with AI Insights")
async def advanced_chat_endpoint(request: ChatRequest):
    """Handle chat messages with advanced AI features"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing advanced message: '{message}'")
        logger.info(f"Jira configured: {app_state.jira_configured}")
        logger.info(f"Jira client exists: {app_state.jira_client is not None}")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            logger.info("🔍 Creating Advanced Chatbot Engine...")
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process query with advanced chatbot
        logger.info("🔍 About to call advanced_chatbot.process_advanced_query")
        result = await app_state.advanced_chatbot.process_advanced_query(message)
        logger.info(f"🔍 Advanced chatbot returned: {result}")
        
        # Generate response
        response = result.get('response', 'I apologize, but I encountered an issue processing your request.')
        logger.info(f"🔍 Final advanced response: {response}")
        
        # Store message in history
        app_state.messages.append({
            "id": len(app_state.messages) + 1,
            "message": message,
            "response": response,
            "timestamp": datetime.now().isoformat(),
            "metadata": {
                "ai_enhanced": True,
                "advanced_features": True,
                "metrics": result.get('metrics', {}),
                "risks": result.get('risks', []),
                "semantic_results": result.get('semantic_results', [])
            }
        })
        
        return {
            "response": response,
            "metadata": {
                "ai_enhanced": True,
                "advanced_features": True,
                "metrics": result.get('metrics', {}),
                "risks": result.get('risks', []),
                "semantic_results": result.get('semantic_results', [])
            }
        }
        
    except Exception as e:
        logger.error(f"Advanced chat error: {e}")
        return {
            "response": f"I encountered an error while processing your request: {str(e)}",
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.get("/api/chat/sprint-health", tags=["Analytics"], summary="Get Sprint Health Dashboard")
async def get_sprint_health():
    """Get comprehensive sprint health analysis"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process sprint health query
        result = await app_state.advanced_chatbot.process_advanced_query("What's our sprint health status?")
        
        return create_success_response({
            "health_dashboard": result.get('response', ''),
            "metrics": result.get('metrics', {}),
            "risks": result.get('risks', [])
        }, "Sprint health analysis completed")
        
    except Exception as e:
        logger.error(f"Sprint health error: {e}")
        return create_error_response("Sprint health analysis failed", str(e))

@app.get("/api/chat/team-performance", tags=["Analytics"], summary="Get Team Performance Analysis")
async def get_team_performance():
    """Get team performance comparison and metrics"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process team performance query
        result = await app_state.advanced_chatbot.process_advanced_query("Compare team performance this sprint")
        
        return create_success_response({
            "team_analysis": result.get('response', ''),
            "comparison_data": result.get('comparison_data', [])
        }, "Team performance analysis completed")
        
    except Exception as e:
        logger.error(f"Team performance error: {e}")
        return create_error_response("Team performance analysis failed", str(e))

@app.post("/api/chat/semantic-search", tags=["Chat"], summary="Semantic Search for Tickets")
async def semantic_search_endpoint(request: ChatRequest):
    """Perform semantic search for tickets"""
    try:
        message = request.message.strip()
        
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process semantic search query
        result = await app_state.advanced_chatbot.process_advanced_query(message)
        
        return create_success_response({
            "search_results": result.get('response', ''),
            "semantic_results": result.get('semantic_results', [])
        }, "Semantic search completed")
        
    except Exception as e:
        logger.error(f"Semantic search error: {e}")
        return create_error_response("Semantic search failed", str(e))

@app.post("/api/jira/test")
async def test_jira_connection(config: JiraConfigRequest):
    """Test Jira connection"""
    try:
        # Create JiraConfig object
        jira_config = JiraConfig(
            base_url=config.url,
            email=config.email,
            api_token=config.api_token,
            board_id=config.board_id
        )
        
        # Create Jira client
        jira_client = JiraClient(jira_config)
        
        # Initialize the async client
        await jira_client.initialize()
        
        # Test the connection
        try:
            # Try to get current sprint
            current_sprint = await jira_client.get_current_sprint()
            sprint_info = f"Current sprint: {current_sprint.get('name', 'Unknown')}" if current_sprint else "No active sprint"
            
            # Try a simple search
            search_result = await jira_client.search("project is not EMPTY", max_results=1)
            total_issues = search_result.get('total', 0)
        except Exception as e:
            logger.warning(f"Could not get current sprint, but connection may still work: {e}")
            sprint_info = "Connection established but sprint info unavailable"
            total_issues = 0
        
        return {
            "success": True,
                "message": f"Jira connection successful! {sprint_info}. Found {total_issues} total issues.",
                "config": {
                    "url": config.url,
                    "email": config.email,
                    "board_id": config.board_id
                },
                "details": {
                    "current_sprint": current_sprint,
                    "total_issues": total_issues
                }
        }
    except Exception as e:
        return {
            "success": False,
            "message": f"Connection test failed: {str(e)}",
            "config": {
                "url": config.url,
                "email": config.email,
                "board_id": config.board_id
            }
        }
    finally:
        # Close the test client
        if 'jira_client' in locals():
            await jira_client.close()

@app.get("/api/leadership/overview")
async def leadership_overview(board_id: int, history: int = 5, future: int = 2):
    """Get comprehensive leadership overview with sprint metrics, velocity, and forecasts"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Ensure Jira client is properly initialized
        if not app_state.jira_client._client:
            await app_state.jira_client.initialize()
        
        # Get current, previous, and future sprints
        sprints = await app_state.jira_client.get_active_sprints()
        closed_sprints = await app_state.jira_client.get_closed_sprints()
        future_sprints = await app_state.jira_client.get_future_sprints()
        
        # Determine current sprint
        current_sprint = None
        previous_sprint = None
        future_sprint = None
        
        if sprints:
            current_sprint = sprints[0]
        if closed_sprints:
            previous_sprint = closed_sprints[-1]  # Most recent closed sprint
        if future_sprints:
            future_sprint = future_sprints[0]
        
        # Get sprint data and calculate metrics
        overview_data = {
            "current_sprint": None,
            "previous_sprint": None,
            "future_sprint": None,
            "velocity_forecast": None,
            "bandwidth_analysis": [],
            "insights": [],
            "summary": {}
        }
        
        # Process current sprint
        if current_sprint:
            sprint_id = current_sprint.get('id')
            sprint_data = await app_state.jira_client.search(f"sprint = {sprint_id}", max_results=1000)
            current_metrics = parse_sprint_metrics(sprint_data.get('issues', []), current_sprint)
            overview_data["current_sprint"] = {
                "id": current_metrics.sprint_id,
                "name": current_metrics.sprint_name,
                "start_date": current_metrics.start_date.isoformat() if current_metrics.start_date else None,
                "end_date": current_metrics.end_date.isoformat() if current_metrics.end_date else None,
                "total_stories": current_metrics.total_stories,
                "completed_stories": current_metrics.completed_stories,
                "in_progress_stories": current_metrics.in_progress_stories,
                "todo_stories": current_metrics.todo_stories,
                "total_story_points": current_metrics.total_story_points,
                "completed_story_points": current_metrics.completed_story_points,
                "velocity": current_metrics.velocity,
                "completion_rate": current_metrics.completion_rate,
                "defect_count": current_metrics.defect_count,
                "defect_ratio": current_metrics.defect_ratio,
                "spillover_count": current_metrics.spillover_count,
                "spillover_rate": current_metrics.spillover_rate
            }
            
            # Analyze bandwidth
            bandwidth = analyze_bandwidth(sprint_data.get('issues', []))
            overview_data["bandwidth_analysis"] = [
                {
                    "team_member": member.team_member,
                    "current_load": member.current_load,
                    "max_capacity": member.max_capacity,
                    "utilization_rate": member.utilization_rate,
                    "available_capacity": member.available_capacity,
                    "recommended_story_points": member.recommended_story_points
                }
                for member in bandwidth
            ]
        
        # Process previous sprint for velocity history
        historical_velocities = []
        if previous_sprint:
            sprint_id = previous_sprint.get('id')
            sprint_data = await app_state.jira_client.search(f"sprint = {sprint_id}", max_results=1000)
            previous_metrics = parse_sprint_metrics(sprint_data.get('issues', []), previous_sprint)
            historical_velocities.append(previous_metrics.velocity)
            
            overview_data["previous_sprint"] = {
                "id": previous_metrics.sprint_id,
                "name": previous_metrics.sprint_name,
                "velocity": previous_metrics.velocity,
                "completion_rate": previous_metrics.completion_rate,
                "defect_ratio": previous_metrics.defect_ratio
            }
        
        # Get more historical data for better forecasting
        if len(historical_velocities) < 3:
            # Get last 3 closed sprints
            recent_closed = closed_sprints[-3:] if len(closed_sprints) >= 3 else closed_sprints
            for sprint in recent_closed:
                sprint_id = sprint.get('id')
                sprint_data = await app_state.jira_client.search(f"sprint = {sprint_id}", max_results=1000)
                metrics = parse_sprint_metrics(sprint_data.get('issues', []), sprint)
                historical_velocities.append(metrics.velocity)
        
        # Generate velocity forecast
                if current_sprint and historical_velocities:
                    current_sprint_data = await app_state.jira_client.search(f"sprint = {current_sprint.get('id')}", max_results=1000)
                    forecast = forecast_velocity_wrapper(historical_velocities, current_sprint_data.get('issues', []))
            overview_data["velocity_forecast"] = {
                "current_velocity": forecast.current_velocity,
                "projected_velocity": forecast.projected_velocity,
                "burn_rate": forecast.burn_rate,
                "days_remaining": forecast.days_remaining,
                "stories_remaining": forecast.stories_remaining,
                "completion_probability": forecast.completion_probability,
                "confidence_level": forecast.confidence_level
            }
        
        # Process future sprint
        if future_sprint:
            overview_data["future_sprint"] = {
                "id": future_sprint.get('id'),
                "name": future_sprint.get('name'),
                "start_date": future_sprint.get('startDate'),
                "end_date": future_sprint.get('endDate')
            }
        
        # Generate insights
        if current_sprint and overview_data["current_sprint"]:
            from metrics_utils import SprintMetrics, VelocityForecast, BandwidthAnalysis
            current_metrics_obj = SprintMetrics(**overview_data["current_sprint"])
            forecast_obj = VelocityForecast(**overview_data["velocity_forecast"]) if overview_data["velocity_forecast"] else None
            bandwidth_obj = [BandwidthAnalysis(**member) for member in overview_data["bandwidth_analysis"]]
            
            insights = generate_insights(current_metrics_obj, forecast_obj, bandwidth_obj)
            overview_data["insights"] = insights
        
        # Generate summary
        overview_data["summary"] = {
            "total_active_sprints": len(sprints),
            "total_closed_sprints": len(closed_sprints),
            "total_future_sprints": len(future_sprints),
            "current_velocity": overview_data["current_sprint"]["velocity"] if overview_data["current_sprint"] else 0,
            "projected_velocity": overview_data["velocity_forecast"]["projected_velocity"] if overview_data["velocity_forecast"] else 0,
            "completion_probability": overview_data["velocity_forecast"]["completion_probability"] if overview_data["velocity_forecast"] else 0,
            "team_members": len(overview_data["bandwidth_analysis"]),
            "insights_count": len(overview_data["insights"])
        }
        
        return create_success_response(overview_data, "Leadership overview retrieved successfully")
        
    except Exception as e:
        logger.error(f"Leadership overview error: {e}")
        return create_error_response("Leadership overview failed", str(e))

@app.get("/api/leadership/overview")
async def leadership_overview(board_id: int, history: int = 5, future: int = 2):
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    jc = app_state.jira_client

    # Sprints
    sprints = await jc.get_sprints_all(board_id, "active,future,closed")
    if not sprints:
        return {"success": True, "data": {"message": "No sprints found for board"}}

    for s in sprints:
        s["start_dt"] = to_dt(s.get("startDate"))
        s["end_dt"] = to_dt(s.get("endDate"))
        s["num"] = parse_sprint_number(s.get("name"))

    active = [s for s in sprints if (s.get("state") or s.get("stateName")) == "active"]
    closed = sorted([s for s in sprints if (s.get("state") or s.get("stateName")) == "closed"],
                    key=lambda x: x.get("end_dt") or datetime.min.replace(tzinfo=timezone.utc))
    future_s = sorted([s for s in sprints if (s.get("state") or s.get("stateName")) == "future"],
                      key=lambda x: x.get("start_dt") or datetime.max.replace(tzinfo=timezone.utc))

    current_sprint = active[0] if active else None
    previous = closed[-history:] if closed else []
    upcoming = future_s[:future] if future_s else []

    fields = await jc.get_fields()
    sp_field_id = find_story_points_field(fields)

    async def metrics_for(s):
        sid = s["id"]
        # Use count() for fast counting, fallback to search().total
        try:
            total_count = await jc.count(f"sprint = {sid}")
        except Exception:
            # Fallback if approximate-count not available
            issues = await jc.search(f"sprint = {sid}", max_results=0)
            total_count = issues.get("total", 0)
        
        # Get actual issues for detailed analysis
        issues = await jc.search(f"sprint = {sid}", max_results=1000)
        items = issues.get("issues", [])
        total_issues = total_count  # Use the grounded count
        done_issues = [i for i in items if is_done(i)]
        total_done = len(done_issues)

        stories = [i for i in items if (i.get("fields", {}).get("issuetype", {}).get("name","").lower() == "story")]
        bugs = [i for i in items if (i.get("fields", {}).get("issuetype", {}).get("name","").lower() in ["bug","defect"])]

        total_points = sum(sum_story_points(i, sp_field_id) for i in items) if sp_field_id else None
        completed_points = sum(sum_story_points(i, sp_field_id) for i in done_issues) if sp_field_id else None
        velocity = completed_points if sp_field_id else total_done

        defect_ratio = (len(bugs) / max(1, len(stories))) * 100
        completion_rate = (total_done / max(1, total_issues)) * 100

        # Bandwidth by assignee (issues basis; add points if needed)
        by_assignee = {}
        for i in items:
            f = i.get("fields", {}) or {}
            assignee = (f.get("assignee", {}) or {}).get("displayName") or "Unassigned"
            by_assignee.setdefault(assignee, {"issues": 0, "done": 0, "points": 0.0})
            by_assignee[assignee]["issues"] += 1
            if is_done(i):
                by_assignee[assignee]["done"] += 1
            if sp_field_id:
                by_assignee[assignee]["points"] += sum_story_points(i, sp_field_id)

        return {
            "sprint_id": sid,
            "sprint_name": s.get("name"),
            "sprint_number": s.get("num"),
            "start": s.get("startDate"),
            "end": s.get("endDate"),
            "total_issues": total_issues,
            "done_issues": total_done,
            "total_points": total_points,
            "completed_points": completed_points,
            "velocity": velocity,
            "completion_rate": round(completion_rate, 2),
            "defect_ratio": round(defect_ratio, 2),
            "bandwidth": by_assignee
        }

    # velocity series (previous + current)
    velocity_series = []
    prev_metrics = []
    for s in previous:
        m = await metrics_for(s)
        prev_metrics.append(m)
        velocity_series.append({"label": m["sprint_name"], "velocity": m["velocity"]})

    current_metrics = None
    if current_sprint:
        current_metrics = await metrics_for(current_sprint)
        velocity_series.append({"label": current_metrics["sprint_name"], "velocity": current_metrics["velocity"]})

    # capacity baseline = avg of last 3 closed sprints
    last3 = prev_metrics[-3:] if len(prev_metrics) >= 3 else prev_metrics
    capacity_baseline = round(sum(m["velocity"] for m in last3)/max(1, len(last3)), 2) if last3 else None

    # projections
    series_vals = [pt["velocity"] for pt in velocity_series][-6:]
    proj = forecast_velocity(series_vals)

    # suggestions (rule-based → no hallucinations)
    suggestions = []
    if proj["avg"] > 0:
        delta_pct = (proj["next"] - proj["avg"]) / proj["avg"] * 100
        if delta_pct >= 10:
            suggestions.append(f"Future velocity likely to rise (~{round(delta_pct,1)}%) compared to recent average.")
        elif delta_pct <= -10:
            suggestions.append(f"Future velocity may drop (~{round(abs(delta_pct),1)}%) vs recent average. Consider scope control or capacity boost.")
        else:
            suggestions.append("Future velocity expected to be stable around recent average.")
    if current_metrics:
        cr = current_metrics["completion_rate"]
        dr = current_metrics["defect_ratio"]
        if cr < 50: suggestions.append("Risk of sprint spillover: less than 50% completed.")
        if cr >= 80: suggestions.append("Sprint completion is healthy (≥80%). Keep momentum.")
        if dr >= 20: suggestions.append("Quality risk: defect ratio ≥ 20%. Increase testing and reviews.")
        if dr <= 5:  suggestions.append("Low defect ratio (≤5%). Quality is excellent.")
        # overload hint
        if current_metrics["bandwidth"]:
            busiest = sorted(current_metrics["bandwidth"].items(), key=lambda kv: kv[1]["issues"], reverse=True)[0]
            name, stats = busiest
            if stats["issues"] >= 8:
                suggestions.append(f"{name} may be overloaded ({stats['issues']} issues). Consider redistribution.")

    payload = {
        "board_id": board_id,
        "timeline": {
            "previous": [{"id": s["id"], "name": s.get("name"), "number": s.get("num")} for s in previous],
            "current": {"id": current_sprint["id"], "name": current_sprint["name"], "number": current_sprint.get("num")} if current_sprint else None,
            "future": [{"id": s["id"], "name": s.get("name"), "number": s.get("num")} for s in upcoming]
        },
        "charts": {
            "velocity_series": velocity_series,
            "completion_rate": current_metrics["completion_rate"] if current_metrics else None,
            "defect_ratio": current_metrics["defect_ratio"] if current_metrics else None,
            "bandwidth": current_metrics["bandwidth"] if current_metrics else {},
            "capacity_baseline": capacity_baseline,
            "projections": {
                "velocity_forecast_next": proj["next"],
                "velocity_forecast_band": proj["band"],
                "velocity_recent_avg": proj["avg"],
                "velocity_trend_slope": proj["trend"]
            }
        },
        "current_sprint": current_metrics,
        "suggestions": suggestions
    }
    return {"success": True, "data": payload}

@app.post("/api/export/pdf", tags=["Export"], summary="Export Chat to PDF")
async def export_pdf():
    """Export chat to PDF using reportlab"""
    try:
        if not REPORTLAB_AVAILABLE:
            raise HTTPException(status_code=500, detail="ReportLab not available. Please install reportlab package.")
        
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        import io
        
        # Create PDF in memory
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=30,
        )
        
        content_style = ParagraphStyle(
            'CustomContent',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=12,
        )
        
        # Build PDF content
        story = []
        story.append(Paragraph("Leadership Quality Tool - Chat Export", title_style))
        story.append(Spacer(1, 20))
        
        # Add chat messages
        for i, message in enumerate(app_state.messages, 1):
            story.append(Paragraph(f"<b>Message {i}:</b> {message.get('message', '')}", content_style))
            story.append(Paragraph(f"<b>Response:</b> {message.get('response', '')}", content_style))
            story.append(Paragraph(f"<b>Timestamp:</b> {message.get('timestamp', '')}", content_style))
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        
        # Get PDF content
        pdf_content = buffer.getvalue()
        buffer.close()
        
        # Store in app state for download
        filename = f"chat_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        app_state.export_files[filename] = pdf_content
        
        return create_success_response({
            "filename": filename,
            "size_bytes": len(pdf_content)
        }, "PDF exported successfully")
        
    except ImportError:
        return create_error_response("PDF export failed", "reportlab package not installed")
    except Exception as e:
        logger.error(f"PDF export error: {e}")
        return create_error_response("PDF export failed", str(e))

@app.post("/api/export/powerpoint", tags=["Export"], summary="Export Chat to PowerPoint")
async def export_powerpoint():
    """Export chat to PowerPoint using python-pptx"""
    try:
        if not PPTX_AVAILABLE:
            raise HTTPException(status_code=500, detail="python-pptx not available. Please install python-pptx package.")
        
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        import io
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Leadership Quality Tool"
        subtitle.text = f"Chat Export - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add content slides
        content_layout = prs.slide_layouts[1]
        
        for i, message in enumerate(app_state.messages, 1):
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Message {i}"
            
            # Format content
            text_frame = content.text_frame
            text_frame.clear()
            
            # Add message
            p = text_frame.paragraphs[0]
            p.text = f"User: {message.get('message', '')}"
            p.font.size = Pt(12)
            
            # Add response
            p = text_frame.add_paragraph()
            p.text = f"Assistant: {message.get('response', '')}"
            p.font.size = Pt(10)
            
            # Add timestamp
            p = text_frame.add_paragraph()
            p.text = f"Time: {message.get('timestamp', '')}"
            p.font.size = Pt(8)
            p.font.italic = True
        
        # Save to memory
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_content = buffer.getvalue()
        buffer.close()
        
        # Store in app state for download
        filename = f"chat_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        app_state.export_files[filename] = pptx_content
        
        return create_success_response({
            "filename": filename,
            "size_bytes": len(pptx_content),
            "slides": len(prs.slides)
        }, "PowerPoint exported successfully")
        
    except ImportError:
        return create_error_response("PowerPoint export failed", "python-pptx package not installed")
    except Exception as e:
        logger.error(f"PowerPoint export error: {e}")
        return create_error_response("PowerPoint export failed", str(e))

@app.get("/api/export/download/{filename}", tags=["Export"], summary="Download Exported File")
async def download_export(filename: str):
    """Download exported file"""
    if filename not in app_state.export_files:
        raise HTTPException(status_code=404, detail="File not found")
    
    content = app_state.export_files[filename]
    
    # Determine content type
    if filename.endswith('.pdf'):
        media_type = 'application/pdf'
    elif filename.endswith('.pptx'):
        media_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    else:
        media_type = 'application/octet-stream'
    
    return StreamingResponse(
        io.BytesIO(content),
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/api/jira/assignee/count", tags=["JIRA"], summary="Count issues by assignee")
async def assignee_count(request: AssigneeCountRequest):
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    try:
        # Resolve assignee to accountId for Jira Cloud
        account_id = await app_state.jira_client.resolve_assignee_to_account_id(request.assignee)
        assignee_clause = f"assignee in ({account_id})" if account_id else f'assignee ~ "{request.assignee}"'

        clauses = [assignee_clause]
        if request.project:
            clauses.append(f"project = {request.project}")
        if request.issuetype:
            clauses.append(f"issuetype = {request.issuetype}")
        if request.status_category:
            clauses.append(f'statusCategory = "{request.status_category}"')
        elif request.not_done:
            clauses.append('statusCategory != Done')

        jql = " AND ".join(clauses)

        # Use fast count with fallback
        try:
            count = await app_state.jira_client.count(jql)
        except Exception:
            res = await app_state.jira_client.search(jql, max_results=0)
            count = int(res.get("total", 0))

        # Friendly message and optional top items
        kind = "assignee_open_tasks" if (request.issuetype or "").lower() == "task" else "assignee_open_stories"
        message = format_friendly_response(kind, count, {"assignee": request.assignee})
        top_items: List[Dict[str, str]] = []
        if count > 0:
            top_items = await fetch_top_items(jql + " ORDER BY created DESC", limit=5)

        return {
            "success": True,
            "jql": jql,
            "count": count,
            "message": message,
            "top": top_items,
            "friendly": (
                message + ("\n\n" + "\n".join([
                    f"- {item['key']} — {item['summary']} ({item['url']})" for item in top_items
                ]) if top_items else "")
            )
        }
    except Exception as e:
        logger.error(f"Assignee count failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)